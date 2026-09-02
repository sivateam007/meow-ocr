#!/usr/bin/env python3
"""
Offline PDF OCR Web App with Auto-Language Detection
Flask application for rendering PDF OCR via web interface
"""

import os
import tempfile
import threading
import uuid
from functools import wraps
import re
import logging
import time
import shutil
import json
import glob
import base64
import io
from flask import Flask, request, render_template, send_file, flash, redirect, url_for, jsonify, session, make_response
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path, pdfinfo_from_path
import pytesseract
from pytesseract import image_to_osd
from PIL import Image, ImageOps, ImageEnhance
import requests
import gc
gc.set_threshold(100, 5, 2)  # More aggressive GC for memory-constrained environments
import asyncio
from concurrent.futures import ThreadPoolExecutor, TimeoutError as _CTimeoutError
if not hasattr(asyncio, 'coroutine'):
    asyncio.coroutine = lambda f: f
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", os.urandom(24).hex())

# Security hardening for session cookies
# Detect production (Render) so Secure cookie is only set over HTTPS.
_IS_PRODUCTION = bool(
    os.environ.get("RENDER_EXTERNAL_URL") or os.environ.get("RENDER_URL")
    or os.environ.get("RENDER_SERVICE_NAME")
)
app.config['SESSION_COOKIE_HTTPONLY'] = True          # not readable by JS (anti-XSS)
app.config['SESSION_COOKIE_SAMESITE'] = "Lax"         # anti-CSRF for cross-site POSTs
app.config['SESSION_COOKIE_SECURE'] = _IS_PRODUCTION  # HTTPS-only in production


# =====================================================================
# Lightweight in-memory rate limiter (per IP) to blunt resource-exhaustion
# / brute-force abuse on the free tier. Entries are pruned lazily.
# NOTE: in-memory only — not a substitute for a CDN/WAF, but free + effective,
# and resets across restarts (which is acceptable for this threat model).
# =====================================================================
_rate_buckets = {}
_rate_lock = threading.Lock()


def _rate_limit(limit, window_seconds):
    """Decorator: allow up to `limit` requests per `window_seconds` per IP.
    Returns HTTP 429 when exceeded."""
    def decorator(fn):
        @wraps(fn)
        def wrapper(*args, **kwargs):
            ip = request.remote_addr or "unknown"
            now = time.time()
            with _rate_lock:
                bucket = [t for t in _rate_buckets.get(ip, []) if now - t < window_seconds]
                bucket.append(now)
                _rate_buckets[ip] = bucket
                exceed = len(bucket) > limit
            if exceed:
                return jsonify({"error": "Too many requests, please slow down."}), 429
            return fn(*args, **kwargs)
        return wrapper
    return decorator


def _check_rate(limit, window_seconds):
    """Inline rate-limit check (True = allowed, False = over limit).
    Returns the same error response the decorator would, or None if allowed."""
    ip = request.remote_addr or "unknown"
    now = time.time()
    with _rate_lock:
        bucket = [t for t in _rate_buckets.get(ip, []) if now - t < window_seconds]
        bucket.append(now)
        _rate_buckets[ip] = bucket
        exceed = len(bucket) > limit
    if exceed:
        return jsonify({"error": "Too many requests, please slow down."}), 429
    return None

# Jinja filter: Unix timestamp to readable date
@app.template_filter('datetimeformat')
def datetimeformat(timestamp):
    import datetime
    return datetime.datetime.fromtimestamp(timestamp).strftime('%Y-%m-%d %H:%M')


@app.context_processor
def inject_globals():
    """Make Firebase config + current user available to every template."""
    import json as _json
    from flask import request as _request
    from ad_config import ads_slots
    from monetag_config import MONETAG, MONETAG_ENABLED
    user = get_user()
    _ads = ads_slots()
    _path = _request.path if _request else "/"
    with progress_lock:
        _doc_count = sum(1 for t in progress_tracker.values() if t.get("status") == "completed")
    return {
        "doc_count": _doc_count,
        "canonical_url": (SITE_URL + _path) if SITE_URL else "/",
        "ga4_id": os.environ.get("GA4_ID", "").strip(),
        "clarity_id": os.environ.get("CLARITY_ID", "").strip(),
        "ads": _ads,
        "ads_enabled": bool(_ads.get("leaderboard") or _ads.get("inpage") or _ads.get("popunder")),
        "monetag_enabled": MONETAG_ENABLED,
        "monetag_sdk": MONETAG.get("vignette", ""),
        "monetag_vignette": MONETAG.get("vignette", ""),
        "monetag_inpage": MONETAG.get("inpage", ""),
        "monetag_popunder": MONETAG.get("popunder", ""),
        "monetag_push": MONETAG.get("push", ""),
        "firebase_api_key": FIREBASE_API_KEY,
        "firebase_auth_domain": FIREBASE_AUTH_DOMAIN,
        "firebase_project_id": FIREBASE_PROJECT_ID,
        "firebase_storage_bucket": FIREBASE_STORAGE_BUCKET,
        "firebase_app_id": FIREBASE_APP_ID,
        "firebase_measurement_id": FIREBASE_MEASUREMENT_ID,
        "firebase_enabled": FIREBASE_ENABLED,
        "current_user_json": _json.dumps(user or {}, ensure_ascii=False),
        "current_user": user,
        "free_docs_without_login": FREE_DOCS_WITHOUT_LOGIN,
    }

# Configuration
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE  # Reject oversized uploads early
ALLOWED_EXTENSIONS = {
    # PDF
    'pdf',
    # Images (OCR needed)
    'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif',
    # Word documents
    'docx', 'doc',
    # PowerPoint
    'pptx', 'ppt',
    # Excel/Spreadsheets
    'xlsx', 'xls', 'csv',
    # Rich Text
    'rtf',
    # OpenDocument formats
    'odt', 'ods', 'odp',
    # HTML
    'html', 'htm',
    # XML/JSON
    'xml', 'json',
    # Plain text
    'txt', 'md'
}
DEFAULT_LANG = 'tam'  # Tamil by default
CHECKPOINT_INTERVAL = 5  # Upload checkpoint to Mega every N pages
BATCH_SIZE = 1  # Pages per batch (lower = less memory per batch)
MEGA_LOGIN_TIMEOUT = 30  # Seconds before Mega login times out
CONVERT_TIMEOUT = 120  # Max seconds for convert_from_path per batch
OCR_TIMEOUT = 300  # Max seconds per page for Tesseract OCR
AUTO_DELETE_DAYS = int(os.environ.get("AUTO_DELETE_DAYS", "2"))  # Cloud auto-delete window (default 2 days)
AUTO_DELETE_SECONDS = AUTO_DELETE_DAYS * 86400

# Anonymous free usage limit (docs convertible without signing in)
FREE_DOCS_WITHOUT_LOGIN = int(os.environ.get("FREE_DOCS_WITHOUT_LOGIN", "1"))
_COOKIE_COUNTER = "scan_docs_done"  # cookie name counting anonymous conversions

# Firebase (Google Sign-in) config — set these env vars to enable sign-in.
FIREBASE_API_KEY = os.environ.get("FIREBASE_API_KEY", "")
FIREBASE_PROJECT_ID = os.environ.get("FIREBASE_PROJECT_ID", "")
FIREBASE_AUTH_DOMAIN = os.environ.get("FIREBASE_AUTH_DOMAIN", "")
FIREBASE_APP_ID = os.environ.get("FIREBASE_APP_ID", "")
FIREBASE_MEASUREMENT_ID = os.environ.get("FIREBASE_MEASUREMENT_ID", "")
FIREBASE_STORAGE_BUCKET = os.environ.get("FIREBASE_STORAGE_BUCKET", "")
FIREBASE_ENABLED = bool(FIREBASE_API_KEY and FIREBASE_PROJECT_ID)

# Progress tracking
progress_lock = threading.Lock()
progress_tracker = {}  # task_id: { ... }

# JSON file persistence for progress_tracker (survives Render restarts)
PROGRESS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "progress_tracker.json")
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ocr-outputs")
_last_save_time = 0
_mega_restore_done = threading.Event()
_mega_restore_error = ""

def _save_progress(force=False):
    """Save progress_tracker to JSON file with throttling (max 1 write/sec)."""
    global _last_save_time
    now = time.time()
    if not force and now - _last_save_time < 2:
        return
    serializable = {}
    with progress_lock:
        for tid, task in progress_tracker.items():
            serializable[tid] = dict(task)
    try:
        with open(PROGRESS_FILE, 'w', encoding='utf-8') as f:
            json.dump(serializable, f, default=str, ensure_ascii=False)
        _last_save_time = now
    except Exception as e:
        logger.error(f"Failed to save progress: {e}")

def _load_progress():
    """Load progress_tracker from JSON file on startup."""
    if not os.path.exists(PROGRESS_FILE):
        return {}
    try:
        with open(PROGRESS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        logger.error(f"Failed to load progress: {e}")
        return {}

def persist_output(task_id):
    """Copy completed output file from temp dir to persistent OUTPUT_DIR."""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return False
        src = task.get("output_path")
        output_filename = task.get("output_filename")
        if not src or not os.path.exists(src) or not output_filename:
            return False
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    dest = os.path.join(OUTPUT_DIR, f"{task_id}_{output_filename}")
    try:
        shutil.copy2(src, dest)
        with progress_lock:
            task["output_path"] = dest
        logger.info(f"Task {task_id}: Output persisted to {dest}")
        return True
    except Exception as e:
        logger.error(f"Task {task_id}: Failed to persist output: {e}")
        return False


def rebuild_completed_from_local():
    """Scan local OUTPUT_DIR and restore completed tasks not in progress_tracker."""
    if not os.path.exists(OUTPUT_DIR):
        return
    now = time.time()
    restored = 0
    for fname in os.listdir(OUTPUT_DIR):
        if not fname.endswith(".txt"):
            continue
        filepath = os.path.join(OUTPUT_DIR, fname)
        if not os.path.isfile(filepath):
            continue
        # Filename format: {task_id}_{output_filename}
        task_id = fname.split("_", 1)[0]
        with progress_lock:
            if task_id in progress_tracker:
                continue
        output_filename = fname[len(task_id) + 1:]
        orig_name = output_filename.rsplit("_ocr.", 1)[0] if "_ocr." in output_filename else output_filename
        with progress_lock:
            progress_tracker[task_id] = {
                "current_page": 0, "status": "completed",
                "result_path": None, "error": None,
                "filename": orig_name, "output_filename": output_filename,
                "output_path": filepath,
                "download_link": None,
                "mega_uploaded": False, "mega_status": "",
                "file_type": "pdf", "detected_language": "",
                "pages_processed": 0, "percentage": 100,
                "download_count": 0, "completed_at": now, "created_at": now
            }
            restored += 1
    if restored:
        logger.info(f"Restored {restored} completed tasks from local {OUTPUT_DIR}")


def _get_memory_mb():
    """Get current RSS memory in MB (Linux only). Returns 0 on other platforms."""
    try:
        if os.name == 'posix':
            with open('/proc/self/status') as f:
                for line in f:
                    if line.startswith('VmRSS:'):
                        return int(line.split()[1]) // 1024
    except Exception:
        pass
    return 0

# Self-keepalive to prevent Render from sleeping during background tasks
_active_tasks = 0
_keepalive_lock = threading.Lock()
_keepalive_thread = None

# Secondary keepalive: each active OCR thread also pings itself
_ocr_keepalive_running = threading.Event()

# Language detection mapping
SCRIPT_TO_LANG = {
    "Tamil": "tam",
    "Latin": "eng",  # Default Latin script to English
    "Devanagari": "hin",  # Hindi (most common Devanagari)
    "Telugu": "tel",
    "Bengali": "ben",
    "Kannada": "kan",
    "Malayalam": "mal",
    "Gujarati": "guj",
    "Punjabi": "pan",
    "Marathi": "mar",
    "Arabic": "ara",
    "Cyrillic": "rus",  # Russian
    "Greek": "ell",
    "Hebrew": "heb",
    "Thai": "tha",
    "Chinese": "chi_sim",  # Simplified Chinese
    "Japanese": "jpn",
    "Korean": "kor",
    "Spanish": "spa",
    "French": "fra",
    "German": "deu",
    "Italian": "ita",
    # Add more as needed
}

# Language code to full name mapping
LANG_CODE_TO_NAME = {
    "tam": "Tamil",
    "eng": "English",
    "hin": "Hindi",
    "tel": "Telugu",
    "ben": "Bengali",
    "kan": "Kannada",
    "mal": "Malayalam",
    "guj": "Gujarati",
    "pan": "Punjabi",
    "mar": "Marathi",
    "ara": "Arabic",
    "rus": "Russian",
    "ell": "Greek",
    "heb": "Hebrew",
    "tha": "Thai",
    "chi_sim": "Chinese (Simplified)",
    "jpn": "Japanese",
    "kor": "Korean",
    "spa": "Spanish",
    "fra": "French",
    "deu": "German",
    "ita": "Italian",
}

# OCR language code to translator language code mapping
TRANSLATOR_LANG_MAP = {
    "tam": "ta",
    "eng": "en",
    "hin": "hi",
    "tel": "te",
    "ben": "bn",
    "kan": "kn",
    "mal": "ml",
    "guj": "gu",
    "pan": "pa",
    "mar": "mr",
    "ara": "ar",
    "rus": "ru",
    "ell": "el",
    "heb": "he",
    "tha": "th",
    "chi_sim": "zh-cn",
    "jpn": "ja",
    "kor": "ko",
    "spa": "es",
    "fra": "fr",
    "deu": "de",
    "ita": "it",
}

# Reverse map: translator code -> display name for translation target languages
TRANSLATOR_TARGET_LANGS = {
    "en": "English",
    "ta": "Tamil",
    "hi": "Hindi",
    "te": "Telugu",
    "bn": "Bengali",
    "kn": "Kannada",
    "ml": "Malayalam",
    "gu": "Gujarati",
    "pa": "Punjabi",
    "mr": "Marathi",
    "ar": "Arabic",
    "ru": "Russian",
    "el": "Greek",
    "he": "Hebrew",
    "th": "Thai",
    "zh-cn": "Chinese (Simplified)",
    "ja": "Japanese",
    "ko": "Korean",
    "es": "Spanish",
    "fr": "French",
    "de": "German",
    "it": "Italian",
    "pt": "Portuguese",
    "tr": "Turkish",
    "vi": "Vietnamese",
    "id": "Indonesian",
    "ms": "Malay",
    "nl": "Dutch",
    "pl": "Polish",
    "sv": "Swedish",
    "da": "Danish",
    "fi": "Finnish",
    "cs": "Czech",
    "ro": "Romanian",
    "uk": "Ukrainian",
    "hu": "Hungarian",
}

def get_file_type(filename):
    """Determine file type and processing method"""
    ext = filename.rsplit('.', 1)[1].lower()
    if ext == 'pdf':
        return 'pdf'
    elif ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif']:
        return 'image'
    elif ext in ['docx', 'doc']:
        return 'docx'
    elif ext in ['pptx', 'ppt']:
        return 'pptx'
    elif ext in ['xlsx', 'xls', 'csv']:
        return 'spreadsheet'
    elif ext == 'rtf':
        return 'rtf'
    elif ext in ['odt', 'ods', 'odp']:
        return 'opendocument'
    elif ext in ['html', 'htm']:
        return 'html'
    elif ext in ['xml', 'json']:
        return 'data'
    elif ext in ['txt', 'md']:
        return 'text'
    return None

def allowed_file(filename):
    """Check if file has allowed extension"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_page_count(pdf_path):
    """
    Get total page count using multiple methods (pdfinfo_from_path, then PyPDF2 fallback)
    """
    # Method 1: pdfinfo_from_path (fast, from poppler)
    try:
        pdf_info = pdfinfo_from_path(pdf_path)
        total_pages = pdf_info.get("Pages", pdf_info.get("pages"))
        logger.info(f"Page count from pdfinfo: {total_pages}")
        return total_pages
    except Exception as e:
        logger.warning(f"pdfinfo failed: {e}, trying PyPDF2")

    # Method 2: PyPDF2 (more reliable fallback)
    if PdfReader:
        try:
            reader = PdfReader(pdf_path)
            total_pages = len(reader.pages)
            logger.info(f"Page count from PyPDF2: {total_pages}")
            return total_pages
        except Exception as e:
            logger.error(f"PyPDF2 failed: {e}")

    logger.warning("Could not determine total page count")
    return None  # Unknown

def detect_pdf_language(pdf_path, sample_page=1):
    """
    Detect the primary language of a PDF by sampling the first page.
    Returns Tesseract language code (e.g., 'tam', 'eng').
    Uses OSD for non-Tamil scripts, verifies with character检查 for Tamil.
    """
    try:
        images = convert_from_path(
            pdf_path,
            dpi=150,
            first_page=sample_page,
            last_page=sample_page
        )
        if not images:
            logger.warning("No image from first page, defaulting to Tamil")
            return DEFAULT_LANG

        # Step 1: OSD (fast, good for non-Tamil scripts)
        osd_result = image_to_osd(images[0])
        logger.info(f"OSD result: {osd_result[:200]}...")

        script_match = re.search(r"Script: ([^\n]+)", osd_result)
        if script_match:
            script_name = script_match.group(1).strip()
            # If OSD says Tamil (or any non-Latin script), trust it
            if script_name != "Latin":
                lang_code = SCRIPT_TO_LANG.get(script_name, DEFAULT_LANG)
                logger.info(f"OSD detected: {script_name} → Language: {lang_code}")
                return lang_code

        # Step 2: OSD said Latin or failed — verify with quick Tamil OCR
        # Many Tamil PDFs are misdetected as Latin by OSD
        try:
            tam_text = pytesseract.image_to_string(images[0], lang='tam')
            tamil_chars = len(re.findall(r'[\u0B80-\u0BFF]', tam_text))
            logger.info(f"Tamil verification: {tamil_chars} Tamil chars found")
            if tamil_chars > 5:
                logger.info("Confirmed Tamil via character detection")
                return 'tam'
        except Exception:
            pass

        if script_match:
            script_name = script_match.group(1).strip()
            lang_code = SCRIPT_TO_LANG.get(script_name, DEFAULT_LANG)
            logger.info(f"No Tamil chars found, using OSD: {script_name} → {lang_code}")
            return lang_code

    except Exception as e:
        logger.warning(f"Language detection failed: {e}, defaulting to Tamil")
    return DEFAULT_LANG


def _ocr_page(img, lang, timeout, task_id, page):
    """Run OCR on a single image with timeout. Uses a fresh thread pool per call.
    The source image is preprocessed (grayscale/upscale/contrast) to boost
    Tesseract accuracy, esp. for Tamil/Indic low-DPI scans."""
    pool = ThreadPoolExecutor(max_workers=1)
    try:
        prepped = preprocess_page(img)
        try:
            f = pool.submit(pytesseract.image_to_string, prepped, lang)
            result = f.result(timeout=timeout)
        finally:
            if prepped is not img:
                prepped.close()
        return result
    except Exception as e:
        logger.warning(f"Task {task_id}: OCR failed on page {page} ({timeout}s timeout): {e}")
        raise
    finally:
        pool.shutdown(wait=False)


def preprocess_page(img):
    """
    Lightweight preprocessing that boosts Tesseract accuracy (esp. Tamil/Indic)
    and trims wasted CPU on overly large scans.

    - Grayscale (Tesseract prefers it).
    - Upscale very small / low-DPI scans toward ~300 DPI so glyphs are readable.
    - Downscale absurdly large images so Tesseract isn't matting millions of pixels.
    - Gentle auto-contrast + slight contrast boost for faint/blurry scans.
    Returns a NEW PIL image (mode 'L'); caller should close it when done.
    """
    try:
        work = img.convert("L")

        w, h = work.size
        long_side = max(w, h)

        # Upscale low-DPI scans toward a reliable ~300 DPI window.
        if long_side < 1400 and min(w, h) > 0:
            scale = round(1400.0 / long_side, 3)
            if 1.0 < scale <= 3.0:
                work = work.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
                w, h = work.size

        # Downscale oversized images to avoid wasting CPU on huge pixel counts.
        if max(w, h) > 2500:
            scale = 2500.0 / max(w, h)
            work = work.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

        try:
            work = ImageOps.autocontrast(work)
            work = ImageEnhance.Contrast(work).enhance(1.15)
        except Exception:
            pass

        return work
    except Exception as e:
        logger.warning(f"Preprocessing failed ({getattr(img, 'mode', '?')}): {e}; using raw image")
        try:
            return img.convert("L")
        except Exception:
            return img


def process_pdf_ocr(pdf_path, lang=DEFAULT_LANG, dpi=200, task_id=None, output_file=None, start_page=None, end_page=None):
    """
    Process PDF pages in batches using OCR, write directly to file for memory efficiency.
    Updates progress_tracker if task_id is provided.
    Processes pages from start_page to end_page (inclusive) in batches of BATCH_SIZE.
    If end_page is None, processes only start_page (single page mode).
    Returns number of pages processed.
    """
    page_num = start_page if start_page else 1
    actual_end = end_page if end_page else page_num
    pages_processed = 0
    batch_size = BATCH_SIZE

    current = page_num
    while current <= actual_end:
        # Pre-batch GC and memory check
        gc.collect()
        mem = _get_memory_mb()
        if mem > 400:
            logger.warning(f"Task {task_id}: RSS {mem}MB > 400MB before page {current}, forcing aggressive GC")
            gc.collect()
            gc.collect()
        
        # Check cancel
        if task_id:
            with progress_lock:
                if progress_tracker[task_id].get("cancelled"):
                    logger.info(f"Task {task_id}: OCR cancelled")
                    return pages_processed
        batch_end = min(current + batch_size - 1, actual_end)

        if task_id:
            with progress_lock:
                progress_tracker[task_id]["current_page"] = current

        convert_pool = ThreadPoolExecutor(max_workers=1)
        try:
            convert_future = convert_pool.submit(
                convert_from_path, pdf_path, dpi=dpi,
                first_page=current, last_page=batch_end
            )
            images = convert_future.result(timeout=CONVERT_TIMEOUT)
        except _CTimeoutError:
            logger.warning(f"Task {task_id}: Batch convert timed out pages {current}-{batch_end}, trying single-page fallback")
            images = None
        except Exception as e:
            logger.warning(f"Task {task_id}: Batch convert failed pages {current}-{batch_end}: {e}, trying single-page fallback")
            images = None
        finally:
            convert_pool.shutdown(wait=False)

        # If batch convert failed, process pages one at a time
        if images is None:
            single_page = current
            while single_page <= batch_end:
                # Check cancel
                if task_id:
                    with progress_lock:
                        if progress_tracker[task_id].get("cancelled"):
                            logger.info(f"Task {task_id}: OCR cancelled during fallback")
                            return pages_processed
                convert_pool2 = ThreadPoolExecutor(max_workers=1)
                try:
                    sf = convert_pool2.submit(
                        convert_from_path, pdf_path, dpi=dpi,
                        first_page=single_page, last_page=single_page
                    )
                    single_images = sf.result(timeout=CONVERT_TIMEOUT)
                except Exception:
                    logger.warning(f"Task {task_id}: Skipping page {single_page} (convert failed/timed out)")
                    single_page += 1
                    continue
                finally:
                    convert_pool2.shutdown(wait=False)

                if single_images:
                    page = single_page
                    t0 = time.time()
                    try:
                        text = _ocr_page(single_images[0], lang, OCR_TIMEOUT, task_id, page)
                        elapsed = time.time() - t0
                        if task_id:
                            with progress_lock:
                                times = progress_tracker[task_id].setdefault("page_times", [])
                                times.append(elapsed)
                                if len(times) > 5:
                                    times.pop(0)
                    except Exception:
                        text = f"[OCR failed on page {page}]"
                    single_images[0].close()
                    if output_file:
                        output_file.write(f"--- Page {page} ---\n{text}\n\n")
                        output_file.flush()
                    pages_processed += 1
                    if task_id:
                        with progress_lock:
                            progress_tracker[task_id]["current_page"] = page
                            total = progress_tracker[task_id].get("total_pages")
                            start_offset = progress_tracker[task_id].get("processing_start_page", 1)
                            if total:
                                relative_page = max(0, page - start_offset + 1)
                                pct = max(1, min(int((relative_page / total) * 100), 99))
                                progress_tracker[task_id]["percentage"] = pct
                    # Explicitly clean up single_images list
                    for sim in single_images:
                        try:
                            sim.close()
                        except Exception:
                            pass
                    single_images.clear()
                gc.collect()
                single_page += 1
            current = batch_end + 1
            gc.collect()
            continue

        if not images:
            break

        # OCR each page with timeout
        for i, img in enumerate(images):
            # Check cancel
            if task_id:
                with progress_lock:
                    if progress_tracker[task_id].get("cancelled"):
                        logger.info(f"Task {task_id}: OCR cancelled mid-batch")
                        img.close()
                        return pages_processed
            page = current + i
            t0 = time.time()
            try:
                text = _ocr_page(img, lang, OCR_TIMEOUT, task_id, page)
                elapsed = time.time() - t0
                if task_id:
                    with progress_lock:
                        times = progress_tracker[task_id].setdefault("page_times", [])
                        times.append(elapsed)
                        if len(times) > 5:
                            times.pop(0)
            except Exception:
                text = f"[OCR failed on page {page}]"
            img.close()
            if output_file:
                output_file.write(f"--- Page {page} ---\n{text}\n\n")
                output_file.flush()
            pages_processed += 1
            if task_id:
                    with progress_lock:
                        progress_tracker[task_id]["current_page"] = page
                        total = progress_tracker[task_id].get("total_pages")
                        start_offset = progress_tracker[task_id].get("processing_start_page", 1)
                        if total:
                            relative_page = max(0, page - start_offset + 1)
                            pct = max(1, min(int((relative_page / total) * 100), 99))
                            progress_tracker[task_id]["percentage"] = pct

        # Aggressive cleanup: close all remaining image handles
        for im in images:
            try:
                im.close()
            except Exception:
                pass
        del images
        gc.collect()
        # Log memory every 20 pages
        if pages_processed > 0 and pages_processed % 20 == 0:
            mem = _get_memory_mb()
            if mem:
                logger.info(f"Task {task_id}: Processed {pages_processed} pages, RSS ~{mem}MB")
        current = batch_end + 1

    return pages_processed

def process_image_file(image_path, lang='eng'):
    """Process image file with OCR"""
    try:
        img = Image.open(image_path)
        try:
            prepped = preprocess_page(img)
        except Exception:
            prepped = img
            if prepped.mode != 'RGB':
                prepped = prepped.convert('RGB')
        try:
            text = pytesseract.image_to_string(prepped, lang=lang)
        finally:
            if prepped is not img:
                prepped.close()
            img.close()
        return text
    except Exception as e:
        logger.error(f"Image processing error: {e}")
        raise


def ocr_low_confidence_words(image_path, lang='eng', threshold=62, max_words=40):
    """
    Run Tesseract with per-word confidence (image OCR only) and return the
    words the engine was unsure about, so we can highlight likely OCR errors.
    Best-effort: any failure returns an empty list and never breaks the job.
    """
    try:
        img = Image.open(image_path)
        try:
            try:
                prepped = preprocess_page(img)
            except Exception:
                prepped = img if img.mode == 'L' else img.convert('L')
            data = pytesseract.image_to_data(prepped, lang=lang, output_type=pytesseract.Output.DICT)
        finally:
            if prepped is not img:
                prepped.close()
            img.close()

        words = []
        for i, w in enumerate(data.get("text", [])):
            try:
                conf = float(data.get("conf", [100])[i])
            except (TypeError, ValueError):
                conf = 100.0
            w = (w or "").strip()
            if not w:
                continue
            if conf >= 0 and conf < threshold:
                words.append((w, int(conf)))
        # Deduplicate preserving order, sort by lowest confidence
        seen = set()
        unique = []
        for w, c in words:
            key = w.lower()
            if key not in seen:
                seen.add(key)
                unique.append((w, c))
        unique.sort(key=lambda t: t[1])
        return [w for w, c in unique[:max_words]]
    except Exception as e:
        logger.warning(f"Confidence scan skipped ({getattr(e, '__class__', e).__name__}): {e}")
        return []

def process_docx_file(docx_path):
    """Extract text from DOCX file (no OCR needed)"""
    try:
        from docx import Document
        doc = Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"DOCX processing error: {e}")
        raise

def process_pptx_file(pptx_path):
    """Extract text from PPTX file (no OCR needed)"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        full_text = []
        for i, slide in enumerate(prs.slides, 1):
            full_text.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text.strip())
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"PPTX processing error: {e}")
        raise

def process_txt_file(txt_path):
    """Read text file directly"""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(txt_path, 'r', encoding='latin-1') as f:
            return f.read()

def process_spreadsheet_file(file_path, ext):
    """Extract text from Excel/CSV files"""
    try:
        import pandas as pd
        full_text = []
        
        if ext in ['xlsx', 'xls']:
            # Read all sheets
            xl_file = pd.ExcelFile(file_path)
            for sheet_name in xl_file.sheet_names:
                full_text.append(f"--- Sheet: {sheet_name} ---")
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                full_text.append(df.to_string(index=False))
        elif ext == 'csv':
            df = pd.read_csv(file_path)
            full_text.append(df.to_string(index=False))
        
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Spreadsheet processing error: {e}")
        raise

def process_rtf_file(rtf_path):
    """Extract text from RTF files"""
    try:
        from pyth import open as pyth_open
        doc = pyth_open(rtf_path)
        return doc.get_text()
    except Exception as e:
        logger.error(f"RTF processing error: {e}")
        raise

def process_opendocument_file(odt_path, ext):
    """Extract text from OpenDocument files"""
    try:
        from odfpy import opendocument
        from xml.etree import ElementTree as ET
        
        doc = opendocument.load(odt_path)
        
        # Extract text from paragraphs
        full_text = []
        for para in doc.getElementsByType(odfpy.text.P):
            text = para.getFirstChildText()
            if text and text.strip():
                full_text.append(text.strip())
        
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"OpenDocument processing error: {e}")
        raise

def process_html_file(html_path):
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logger.error(f"HTML processing error: {e}")
        raise

def process_data_file(data_path, ext):
    """Extract text from XML/JSON files"""
    try:
        with open(data_path, 'r', encoding='utf-8') as f:
            if ext == 'json':
                import json
                data = json.load(f)
                return json.dumps(data, indent=2)
            else:  # XML
                from xml.etree import ElementTree as ET
                tree = ET.parse(f)
                root = tree.getroot()
                return ET.tostring(root, encoding='unicode', method='text')
    except Exception as e:
        logger.error(f"Data file processing error: {e}")
        raise

def process_file_background(task_id, file_path, filename, temp_dir, selected_lang='auto', start_page=1, end_page=None, dpi=200, file_type='pdf'):
    """
    Background thread function to process any file type and update progress
    """
    try:
        # Initialize progress
        with progress_lock:
            progress_tracker[task_id]["status"] = "processing"
            progress_tracker[task_id]["total_pages"] = 1  # Single unit for non-PDF
        
        # Create output file path
        output_filename = f"{os.path.splitext(filename)[0]}_ocr.txt"
        output_path = os.path.join(temp_dir, output_filename)
        
        with progress_lock:
            progress_tracker[task_id]["output_path"] = output_path
            progress_tracker[task_id]["output_filename"] = output_filename
        
        # Process based on file type
        if file_type == 'pdf':
            # Existing PDF processing logic
            # Determine language
            if selected_lang != 'auto':
                # Manual selection
                detected_lang = selected_lang
                logger.info(f"Task {task_id}: Using manually selected language {detected_lang}")
                with progress_lock:
                    progress_tracker[task_id]["status"] = "getting_page_count"
                    progress_tracker[task_id]["detected_language"] = detected_lang
            else:
                # Auto-detect
                with progress_lock:
                    progress_tracker[task_id]["status"] = "detecting_language"
                
                detected_lang = detect_pdf_language(file_path)
                logger.info(f"Task {task_id}: Detected language {detected_lang}")
                
                with progress_lock:
                    progress_tracker[task_id]["status"] = "getting_page_count"
                    progress_tracker[task_id]["detected_language"] = detected_lang
            
            logger.info(f"Task {task_id}: Starting page count for {file_path}")
            
            # Get total page count
            total_pages = get_page_count(file_path)
            logger.info(f"Task {task_id}: Total pages = {total_pages}")
            
            # Adjust page range
            actual_start = start_page
            actual_end = end_page
            
            if actual_end is None or actual_end > total_pages:
                actual_end = total_pages
            
            if actual_start > total_pages:
                actual_start = 1
            
            logger.info(f"Task {task_id}: Processing pages {actual_start} to {actual_end}")
            
            with progress_lock:
                progress_tracker[task_id]["status"] = "processing"
                page_range = actual_end - actual_start + 1 if actual_end else total_pages - actual_start + 1
                progress_tracker[task_id]["total_pages"] = page_range
                progress_tracker[task_id]["pdf_total_pages"] = total_pages
                progress_tracker[task_id]["processing_start_page"] = actual_start
            
            # Process PDF and write directly to file (memory efficient)
            logger.info(f"Task {task_id}: Starting OCR processing with language {detected_lang}")
            
            mega_ckpt = None
            originals_uploaded = False
            
            # Track start time
            ocr_start_time = time.time()
            pages_processed = 0
            last_checkpoint_pages = 0
            
            with open(output_path, 'w', encoding='utf-8') as output_file:
                current = actual_start
                while actual_end is None or current <= actual_end:
                    if total_pages and current > total_pages:
                        break
                    
                    # Check for cancel
                    cancelled_flag = False
                    if task_id:
                        with progress_lock:
                            if progress_tracker[task_id].get("cancelled"):
                                logger.info(f"Task {task_id}: Cancelled by user")
                                progress_tracker[task_id]["status"] = "cancelled"
                                progress_tracker[task_id]["error"] = "Cancelled by user"
                                cancelled_flag = True
                    if cancelled_flag:
                        _save_progress(True)
                        return
                    
                    batch_end = min(current + BATCH_SIZE - 1, actual_end) if actual_end else current + BATCH_SIZE - 1
                    
                    # Process batch of pages
                    result = process_pdf_ocr(
                        file_path,
                        lang=detected_lang,
                        dpi=dpi,
                        task_id=task_id,
                        output_file=output_file,
                        start_page=current,
                        end_page=batch_end
                    )
                    if result == 0:
                        break
                    
                    pages_processed += result
                    current = batch_end + 1
                    
                    # Log page progress every 10 pages
                    if pages_processed % 10 == 0:
                        logger.info(f"Task {task_id}: Processed {pages_processed} pages so far")
                        _ocr_keepalive_ping()
                    
                    # Save progress every page for crash recovery (throttled to max 1 write/2s)
                    _save_progress()
                    
                    # Upload checkpoint to Mega every CHECKPOINT_INTERVAL pages
                    if pages_processed - last_checkpoint_pages >= CHECKPOINT_INTERVAL:
                        last_checkpoint_pages = pages_processed
                        # Lazy Mega login (only when first checkpoint is due)
                        if mega_ckpt is None and os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD"):
                            mega_ckpt = init_mega()

                        if mega_ckpt:
                            # Upload original file on first checkpoint only (best-effort)
                            if not originals_uploaded:
                                try:
                                    originals_handle = ensure_mega_folder(mega_ckpt, "ocr-originals")
                                    if originals_handle:
                                        result = mega_call(mega_ckpt, "upload", file_path, dest=originals_handle, dest_filename=filename, timeout=120)
                                        with progress_lock:
                                            progress_tracker[task_id]["mega_original_handle"] = str(getattr(result, 'node_id', result))
                                        originals_uploaded = True
                                        logger.info(f"Task {task_id}: Original uploaded to Mega for resume")
                                except Exception as e:
                                    logger.warning(f"Task {task_id}: Original upload failed (checkpoint still saved): {e}")

                            checkpoint_data = {
                                "task_id": task_id,
                                "last_page": current - 1,
                                "total_pages": total_pages,
                                "filename": filename,
                                "output_filename": output_filename,
                                "detected_lang": detected_lang,
                                "file_type": "pdf",
                                "start_page": actual_start,
                                "end_page": actual_end,
                                "created_at": time.time(),
                                "original_filename": filename
                            }
                            try:
                                upload_checkpoint(mega_ckpt, task_id, output_path, checkpoint_data)
                                logger.info(f"Task {task_id}: Checkpoint saved at page {current - 1}")
                                with progress_lock:
                                    progress_tracker[task_id]["last_checkpoint_page"] = current - 1
                                _save_progress()
                            except Exception as e:
                                logger.warning(f"Task {task_id}: Checkpoint upload failed: {e}")
            
            logger.info(f"Task {task_id}: OCR returned {pages_processed} pages")
            
            # Check if cancelled before treating empty output as error
            if pages_processed == 0:
                cancelled_here = False
                with progress_lock:
                    if progress_tracker[task_id].get("cancelled"):
                        logger.info(f"Task {task_id}: Cancelled by user")
                        progress_tracker[task_id]["status"] = "cancelled"
                        progress_tracker[task_id]["error"] = "Cancelled by user"
                        cancelled_here = True
                if cancelled_here:
                    _save_progress(True)
                    return

            # Verify output file has content
            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                logger.error(f"Task {task_id}: Output file empty or missing")
                with progress_lock:
                    progress_tracker[task_id]["status"] = "error"
                    progress_tracker[task_id]["error"] = "No text could be extracted"
                _save_progress(True)
                return
            
            # Update tracker with success (check cancel race)
            with progress_lock:
                if progress_tracker[task_id].get("cancelled"):
                    progress_tracker[task_id]["status"] = "cancelled"
                    progress_tracker[task_id]["error"] = "Cancelled by user"
                else:
                    progress_tracker[task_id]["status"] = "completed"
                    progress_tracker[task_id]["pages_processed"] = pages_processed
                    progress_tracker[task_id]["percentage"] = 100
            _save_progress(True)
            logger.info(f"Task {task_id}: OCR completed successfully")
            
        elif file_type == 'image':
            # Update language detection for images
            if selected_lang == 'auto':
                with progress_lock:
                    progress_tracker[task_id]["detected_language"] = 'eng'  # Default for images
            
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            # Process image with OCR
            text = process_image_file(file_path, lang=selected_lang if selected_lang != 'auto' else 'eng')
            
            # Confidence highlighting (image OCR only) — best-effort, never breaks the job
            low_conf = ocr_low_confidence_words(file_path, lang=selected_lang if selected_lang != 'auto' else 'eng')
            with progress_lock:
                progress_tracker[task_id]["low_conf_words"] = low_conf
            
            # Write to output
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'docx':
            # Extract text directly (no OCR)
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_docx_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'pptx':
            # Extract text from slides
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_pptx_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'txt':
            # Read directly
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_txt_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'spreadsheet':
            # Extract text from spreadsheets
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            ext = filename.rsplit('.', 1)[1].lower()
            text = process_spreadsheet_file(file_path, ext)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'rtf':
            # Extract text from RTF
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_rtf_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'opendocument':
            # Extract text from OpenDocument
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            ext = filename.rsplit('.', 1)[1].lower()
            text = process_opendocument_file(file_path, ext)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'html':
            # Extract text from HTML
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            text = process_html_file(file_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
                
        elif file_type == 'data':
            # Extract text from XML/JSON
            with progress_lock:
                progress_tracker[task_id]["current_page"] = 1
                progress_tracker[task_id]["percentage"] = 50
            
            ext = filename.rsplit('.', 1)[1].lower()
            text = process_data_file(file_path, ext)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            with progress_lock:
                progress_tracker[task_id]["percentage"] = 100
                progress_tracker[task_id]["pages_processed"] = 1
        
        # Save progress for non-PDF files before Mega upload attempt
        _save_progress()
        
        # Upload to Mega.nz cloud for permanent storage
        try:
            mega_link = upload_to_mega(output_path, output_filename)
            with progress_lock:
                progress_tracker[task_id]["download_link"] = mega_link
                if mega_link:
                    progress_tracker[task_id]["mega_uploaded"] = True
                    progress_tracker[task_id]["mega_status"] = "uploaded"
                    logger.info(f"Task {task_id}: Mega upload success - {mega_link}")
                    # Clean up checkpoint files on successful upload
                    try:
                        m_clean = init_mega()
                        if m_clean:
                            cleanup_checkpoints(m_clean, task_id)
                    except Exception:
                        pass
                else:
                    progress_tracker[task_id]["mega_uploaded"] = False
                    progress_tracker[task_id]["mega_status"] = ""
                    logger.warning(f"Task {task_id}: Mega upload returned None (check Render logs)")
        except Exception as mega_err:
            error_msg = str(mega_err)
            logger.error(f"Task {task_id}: Mega upload error: {error_msg}")
            with progress_lock:
                progress_tracker[task_id]["mega_uploaded"] = False
                progress_tracker[task_id]["mega_status"] = ""

        # Persist output file to ocr-outputs folder
        persist_output(task_id)

        # Mark as completed (respect cancel flag)
        with progress_lock:
            if not progress_tracker[task_id].get("cancelled"):
                progress_tracker[task_id]["status"] = "completed"
            progress_tracker[task_id]["completed_at"] = time.time()
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0)
        _save_progress(True)
            
    except Exception as e:
        logger.error(f"Task {task_id}: Error - {str(e)}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")
        with progress_lock:
            progress_tracker[task_id]["status"] = "error"
            progress_tracker[task_id]["error"] = str(e)
        _save_progress(True)
    finally:
        _release_keepalive()

def cleanup_old_tasks():
    """Clean up completed tasks older than AUTO_DELETE_DAYS (default 2 days).
    Deletes local files and Mega cloud files (auto-delete).
    Preserves interrupted/error tasks so users can retry."""
    current_time = time.time()
    to_delete = []
    mega_delete_names = []
    with progress_lock:
        for task_id, task in progress_tracker.items():
            status = task.get("status", "")
            # Never auto-delete interrupted, error, cancelled, or processing tasks
            if status in ("interrupted", "error", "cancelled", "cancelling", "processing", "resuming", "starting", "detecting_language", "getting_page_count"):
                continue
            if "created_at" not in task:
                continue
            # Respect per-task retention; 0 = keep forever
            retention_days = task.get("auto_delete_days", AUTO_DELETE_DAYS)
            try:
                retention_days = int(retention_days)
            except (TypeError, ValueError):
                retention_days = AUTO_DELETE_DAYS
            if retention_days == 0:
                continue
            retention_seconds = retention_days * 86400
            if current_time - task["created_at"] > retention_seconds:
                to_delete.append(task_id)
                fn = task.get("output_filename")
                if fn:
                    mega_delete_names.append(fn)

        for task_id in to_delete:
            logger.info(f"Auto-deleting old completed task {task_id}")
            op = progress_tracker[task_id].get("output_path")
            if op and os.path.exists(op):
                try: os.remove(op)
                except Exception: pass
            # Remove persisted file from OUTPUT_DIR too
            out_fn = progress_tracker[task_id].get("output_filename")
            if out_fn:
                persisted = os.path.join(OUTPUT_DIR, f"{task_id}_{out_fn}")
                if os.path.exists(persisted):
                    try: os.remove(persisted)
                    except Exception: pass
            if progress_tracker[task_id].get("temp_dir"):
                shutil.rmtree(progress_tracker[task_id]["temp_dir"], ignore_errors=True)
            del progress_tracker[task_id]
    if to_delete:
        _save_progress()
        # Also remove files from Mega cloud storage (auto-delete)
        if mega_delete_names and os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD"):
            mega_delete_files(mega_delete_names)


def mega_delete_files(names, timeout=30):
    """Delete given file names from Mega ocr-outputs folder (cloud auto-delete)."""
    m = init_mega()
    if not m:
        logger.warning("Mega not available for auto-delete")
        return
    try:
        folder = mega_call(m, "find", "ocr-outputs", timeout=timeout)
        if isinstance(folder, (list, tuple)):
            folder = folder[0] if folder else None
        if not folder:
            return
        files = mega_call(m, "get_files_in_node", folder, timeout=timeout)
        if not files:
            return
        for nid, finfo in files.items():
            if not isinstance(finfo, dict):
                continue
            name = finfo.get('a', {}).get('n', '')
            if name in names:
                logger.info(f"Mega auto-delete: removing {name}")
                try:
                    mega_call(m, "delete", nid, timeout=timeout)
                except Exception as e:
                    logger.warning(f"Mega auto-delete failed for {name}: {e}")
    except Exception as e:
        logger.error(f"Mega auto-delete scan error: {e}")


def upload_to_mega(local_file_path, remote_filename):
    """Upload a file to Mega.nz ocr-outputs folder and return download link"""
    if not os.path.exists(local_file_path):
        logger.error(f"File not found for upload: {local_file_path}")
        return None

    m = init_mega()
    if not m:
        logger.warning("Mega login failed - skipping cloud upload")
        return None

    try:
        folders = mega_call(m, "find", "ocr-outputs")
    except Exception:
        logger.warning("Mega find failed during upload")
        return None
    if not folders:
        logger.info("Creating ocr-outputs folder in Mega")
        try:
            folder_node = mega_call(m, "create_folder", "ocr-outputs")
        except Exception:
            logger.warning("Mega create_folder failed during upload")
            return None
        dest = folder_node.get("ocr-outputs")
        if not dest:
            logger.error(f"Failed to get handle from create_folder result: {folder_node}")
            return None
    else:
        dest = folders[0] if isinstance(folders, (list, tuple)) else folders

    logger.info(f"Uploading {remote_filename} to Mega.nz...")
    logger.info(f"File size: {os.path.getsize(local_file_path)} bytes")
    try:
        file_node = mega_call(m, "upload", local_file_path, dest=dest, dest_filename=remote_filename, timeout=120)
    except Exception:
        logger.warning("Mega upload failed")
        return None
    logger.info(f"Upload completed, getting link...")
    try:
        link = mega_call(m, "get_upload_link", file_node)
    except Exception:
        logger.warning("Mega get_upload_link failed")
        return None
    logger.info(f"Mega.nz upload complete: {link}")
    return link


def rebuild_completed_from_mega():
    """Scan Mega ocr-outputs folder and restore completed tasks from _ocr.txt files"""
    global _mega_restore_error
    email = os.environ.get("MEGA_EMAIL")
    password = os.environ.get("MEGA_PWD")
    if not email or not password:
        logger.info("MEGA_EMAIL/MEGA_PWD not set — skipping Mega restore scan")
        return

    try:
        m = init_mega(retries=2)
        if not m:
            _mega_restore_error = "Mega login failed during cloud restore scan"
            logger.error(_mega_restore_error)
            return

        folder = mega_call(m, "find", "ocr-outputs", timeout=15)
        if isinstance(folder, (list, tuple)):
            folder = folder[0] if folder else None
        if not folder:
            logger.info("No ocr-outputs folder found in Mega — nothing to restore")
            return

        try:
            files = mega_call(m, "get_files_in_node", folder, timeout=15)
        except Exception as e:
            logger.warning(f"Mega get_files_in_node failed: {e}")
            return

        if not files:
            logger.info("Mega ocr-outputs folder is empty — nothing to restore")
            return

        now = time.time()
        restored = 0
        for nid, finfo in files.items():
            if not isinstance(finfo, dict):
                continue
            name = finfo.get('a', {}).get('n', '')
            if name.startswith('_'):
                continue
            is_ocr = name.endswith('_ocr.txt')
            is_translation = '_translated_' in name and name.endswith('.txt')
            if not is_ocr and not is_translation:
                continue

            hashlib = __import__('hashlib')
            tid = "mega_" + hashlib.md5(name.encode()).hexdigest()[:12]

            with progress_lock:
                if tid in progress_tracker:
                    continue

            if is_ocr:
                orig_name = name[:-8].rstrip('_')
                file_type = "pdf"
            else:
                orig_name = name
                file_type = "translation"
            with progress_lock:
                progress_tracker[tid] = {
                    "current_page": 0, "status": "completed",
                    "result_path": None, "error": None,
                    "filename": orig_name, "output_filename": name,
                    "download_link": None,
                    "mega_node_id": nid,
                    "mega_node_info": finfo,
                    "mega_uploaded": True, "mega_status": "uploaded",
                    "file_type": file_type, "detected_language": "",
                    "pages_processed": 0, "percentage": 100,
                    "download_count": 0, "completed_at": now, "created_at": now
                }
                restored += 1

        _mega_restore_error = ""
        logger.info(f"Mega restore scan complete: {restored} tasks restored")
    except Exception as e:
        _mega_restore_error = f"Failed to scan Mega for completed tasks: {e}"
        logger.error(_mega_restore_error)


def mega_call(m, method_name, *args, timeout=MEGA_LOGIN_TIMEOUT, **kwargs):
    """Call a Mega client method with timeout protection. Returns the method's result or raises."""
    fn = getattr(m, method_name)
    pool = ThreadPoolExecutor(max_workers=1)
    try:
        future = pool.submit(fn, *args, **kwargs)
        return future.result(timeout=timeout)
    except _CTimeoutError:
        logger.warning(f"Mega.{method_name} timed out after {timeout}s")
        raise
    finally:
        pool.shutdown(wait=False)


def init_mega(retries=1, backoff=3.0):
    """Initialize and login to Mega.nz with timeout. Returns client or None.

    Retries a few times with a short backoff to handle transient/datacenter
    flakiness (e.g. Render's shared outbound IPs occasionally timing out on
    first attempt). Pass retries=3 for startup restore paths.
    """
    email = os.environ.get("MEGA_EMAIL")
    password = os.environ.get("MEGA_PWD")
    if not email or not password:
        return None
    last_err = None
    for attempt in range(retries + 1):
        try:
            from mega import Mega
            from concurrent.futures import ThreadPoolExecutor, TimeoutError as _TimeoutError
            mega = Mega()
            ex = ThreadPoolExecutor(max_workers=1)
            try:
                future = ex.submit(mega.login, email, password)
                return future.result(timeout=MEGA_LOGIN_TIMEOUT)
            finally:
                ex.shutdown(wait=False)
        except _TimeoutError:
            last_err = "Mega login timed out (network issue?)"
            logger.warning(f"{last_err} (attempt {attempt + 1}/{retries + 1})")
        except Exception as e:
            last_err = f"Mega login failed: {e}"
            logger.error(f"{last_err} (attempt {attempt + 1}/{retries + 1})")
            # Don't retry on auth errors (wrong creds) — pointless & slow
            if any(k in str(e).lower() for k in ("bad", "invalid", "authentication", "unauthorized", "wrong", "captcha", "blocked")):
                break
        if attempt < retries:
            time.sleep(backoff)
    return None


def ensure_mega_folder(m, folder_name):
    """Find Mega folder, create if needed. Return folder handle string."""
    try:
        folder = mega_call(m, "find", folder_name)
    except Exception:
        logger.warning(f"Mega find failed for folder {folder_name}")
        return None
    if folder:
        return folder[0]
    try:
        result = mega_call(m, "create_folder", folder_name)
    except Exception:
        logger.warning(f"Mega create_folder failed for {folder_name}")
        return None
    return result.get(folder_name)


def upload_checkpoint(m, task_id, output_path, metadata):
    """Upload checkpoint (partial output + metadata) to ocr-checkpoints/ folder in Mega."""
    import json, tempfile as _tf
    folder_name = "ocr-checkpoints"
    folder_handle = ensure_mega_folder(m, folder_name)
    if not folder_handle:
        logger.error(f"Cannot access/create {folder_name} in Mega")
        return False

    # Delete old checkpoint file if exists
    try:
        old_ckpt = mega_call(m, "find", f"{folder_name}/{task_id}.checkpoint")
        if old_ckpt:
            mega_call(m, "delete", old_ckpt[0])
    except Exception:
        pass
    # Delete old output file if exists
    try:
        old_out = mega_call(m, "find", f"{folder_name}/{task_id}_output.txt")
        if old_out:
            mega_call(m, "delete", old_out[0])
    except Exception:
        pass

    # Upload metadata JSON first
    ckpt_file = _tf.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
    json.dump(metadata, ckpt_file)
    ckpt_file.close()
    try:
        mega_call(m, "upload", ckpt_file.name, dest=folder_handle, dest_filename=f"{task_id}.checkpoint", timeout=120)
    except Exception as e:
        logger.error(f"Checkpoint metadata upload failed: {e}")
        os.unlink(ckpt_file.name)
        return False
    os.unlink(ckpt_file.name)

    # Then upload partial output file
    try:
        mega_call(m, "upload", output_path, dest=folder_handle, dest_filename=f"{task_id}_output.txt", timeout=120)
    except Exception as e:
        logger.error(f"Checkpoint output upload failed: {e}")
        return False

    logger.info(f"Checkpoint saved for task {task_id} at page {metadata.get('last_page')}")
    return True


def cleanup_checkpoints(m, task_id):
    """Delete checkpoint files for a completed task."""
    folder_name = "ocr-checkpoints"
    try:
        old_ckpt = mega_call(m, "find", f"{folder_name}/{task_id}.checkpoint")
        if old_ckpt:
            mega_call(m, "delete", old_ckpt[0])
    except Exception:
        pass
    try:
        old_out = mega_call(m, "find", f"{folder_name}/{task_id}_output.txt")
        if old_out:
            mega_call(m, "delete", old_out[0])
    except Exception:
        pass
    logger.info(f"Checkpoint files cleaned up for task {task_id}")


def scan_and_resume_checkpoints():
    """Called at app startup. Scans Mega for checkpoint files and resumes tasks."""
    logger.info("Scanning for incomplete tasks to resume...")
    try:
        m = init_mega()
        if not m:
            logger.info("Mega not configured - skipping resume scan")
            return

        folder_handle = ensure_mega_folder(m, "ocr-checkpoints")
        if not folder_handle:
            logger.info("No ocr-checkpoints folder - nothing to resume")
            return

        files_in_folder = mega_call(m, "get_files_in_node", folder_handle)
        if not files_in_folder:
            logger.info("No checkpoint files found")
            return

        import json, tempfile as _tf

        for handle, node in files_in_folder.items():
            name = node.get('a', {}).get('n', '')
            if not name.endswith('.checkpoint'):
                continue
            task_id = name[:-len('.checkpoint')]

            logger.info(f"Found incomplete task {task_id}, attempting resume...")
            try:
                temp_dir = _tf.mkdtemp()
                mega_call(m, "download", (handle, node), dest_path=temp_dir)

                ckpt_path = os.path.join(temp_dir, name)
                with open(ckpt_path, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)

                output_name = f"{task_id}_output.txt"
                output_node = mega_call(m, "find", f"ocr-checkpoints/{output_name}")
                if not output_node:
                    logger.warning(f"Task {task_id}: Partial output not found, skipping")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    continue
                mega_call(m, "download", output_node, dest_path=temp_dir)

                original_name = metadata.get("original_filename", "")
                original_dl_path = None
                if original_name:
                    orig_node = mega_call(m, "find", f"ocr-originals/{original_name}")
                    if orig_node:
                        mega_call(m, "download", orig_node, dest_path=temp_dir)
                        original_dl_path = os.path.join(temp_dir, original_name)

                last_page = metadata.get("last_page", 0)
                total_pages = metadata.get("total_pages", 0)
                resume_from = last_page + 1

                if not original_dl_path:
                    logger.warning(f"Task {task_id}: Original file not found, cannot resume")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    continue

                if resume_from > total_pages:
                    logger.info(f"Task {task_id}: Already complete (page {last_page}/{total_pages})")
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    continue

                output_filename = metadata.get("output_filename", f"{os.path.splitext(original_name)[0]}_ocr.txt")
                partial_path = os.path.join(temp_dir, output_name)
                new_output_path = os.path.join(temp_dir, output_filename)

                if os.path.exists(partial_path):
                    os.rename(partial_path, new_output_path)

                detected_lang = metadata.get("detected_lang", "tam")
                file_type = metadata.get("file_type", "pdf")
                actual_start = metadata.get("start_page", 1)
                actual_end = metadata.get("end_page", total_pages)

                with progress_lock:
                    if task_id in progress_tracker:
                        logger.info(f"Task {task_id} already in tracker, skipping")
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        continue
                    progress_tracker[task_id] = {
                        "current_page": resume_from,
                        "status": "resuming",
                        "result_path": None,
                        "error": None,
                        "filename": metadata.get("filename", original_name),
                        "temp_dir": temp_dir,
                        "total_pages": total_pages,
                        "percentage": int((last_page / total_pages) * 100) if total_pages else 0,
                        "detected_language": detected_lang,
                        "created_at": time.time(),
                        "file_type": file_type,
                        "output_path": new_output_path,
                        "output_filename": output_filename,
                        "resumed_from_page": resume_from
                    }

                logger.info(f"Resuming task {task_id} from page {resume_from}/{total_pages}")
                thread = threading.Thread(
                    target=resume_ocr_processing,
                    args=(task_id, original_dl_path, metadata, new_output_path, temp_dir)
                )
                thread.daemon = True
                thread.start()
                _ensure_keepalive()

            except Exception as e:
                logger.error(f"Failed to resume task {task_id}: {e}")
                import traceback
                logger.error(traceback.format_exc())

    except Exception as e:
        logger.error(f"Checkpoint scan error: {e}")


def resume_ocr_processing(task_id, original_path, metadata, output_path, temp_dir):
    """Resume OCR from last checkpoint page."""
    try:
        last_page = metadata.get("last_page", 0)
        total_pages = metadata.get("total_pages", 0)
        detected_lang = metadata.get("detected_lang", "tam")
        actual_start = metadata.get("start_page", 1)
        actual_end = metadata.get("end_page", total_pages)
        filename = metadata.get("filename", "document.pdf")
        output_filename = metadata.get("output_filename", "document_ocr.txt")

        resume_from = last_page + 1
        if resume_from > actual_end or (total_pages and resume_from > total_pages):
            with progress_lock:
                progress_tracker[task_id]["status"] = "completed"
                progress_tracker[task_id]["percentage"] = 100
            logger.info(f"Task {task_id}: Already complete, nothing to resume")
            return

        with progress_lock:
            progress_tracker[task_id]["status"] = "processing"
            progress_tracker[task_id]["total_pages"] = total_pages
        mega_ckpt = init_mega() if os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD") else None

        pages_processed = last_page
        last_checkpoint_pages = last_page

        with open(output_path, 'a', encoding='utf-8') as output_file:
            current = resume_from
            while actual_end is None or current <= actual_end:
                if total_pages and current > total_pages:
                    break

                batch_end = min(current + BATCH_SIZE - 1, actual_end) if actual_end else current + BATCH_SIZE - 1

                result = process_pdf_ocr(
                    original_path,
                    lang=detected_lang,
                    dpi=200,
                    task_id=task_id,
                    output_file=output_file,
                    start_page=current,
                    end_page=batch_end
                )
                if result == 0:
                    break

                pages_processed += result
                current = batch_end + 1

                if mega_ckpt and (pages_processed - last_checkpoint_pages >= CHECKPOINT_INTERVAL):
                    last_checkpoint_pages = pages_processed
                    try:
                        checkpoint_data = {
                            "task_id": task_id,
                            "last_page": current - 1,
                            "total_pages": total_pages,
                            "filename": filename,
                            "output_filename": output_filename,
                            "detected_lang": detected_lang,
                            "file_type": "pdf",
                            "start_page": actual_start,
                            "end_page": actual_end,
                            "created_at": time.time(),
                            "original_filename": filename
                        }
                        upload_checkpoint(mega_ckpt, task_id, output_path, checkpoint_data)
                    except Exception as e:
                        logger.warning(f"Task {task_id}: Resume checkpoint failed: {e}")

        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            logger.error(f"Task {task_id}: Resumed output empty")
            with progress_lock:
                progress_tracker[task_id]["status"] = "error"
                progress_tracker[task_id]["error"] = "Resumed OCR produced no output"
            return

        try:
            mega_link = upload_to_mega(output_path, output_filename)
            with progress_lock:
                progress_tracker[task_id]["download_link"] = mega_link
                if mega_link:
                    progress_tracker[task_id]["mega_uploaded"] = True
                    progress_tracker[task_id]["mega_status"] = "uploaded"
                    if mega_ckpt:
                        try:
                            cleanup_checkpoints(mega_ckpt, task_id)
                        except Exception:
                            pass
                else:
                    progress_tracker[task_id]["mega_uploaded"] = False
                    progress_tracker[task_id]["mega_status"] = ""
        except Exception as mega_err:
            logger.error(f"Task {task_id}: Resume final upload error: {mega_err}")
            with progress_lock:
                progress_tracker[task_id]["mega_uploaded"] = False
                progress_tracker[task_id]["mega_status"] = ""

        persist_output(task_id)

        with progress_lock:
            progress_tracker[task_id]["status"] = "completed"
            progress_tracker[task_id]["pages_processed"] = pages_processed
            progress_tracker[task_id]["percentage"] = 100
            progress_tracker[task_id]["completed_at"] = time.time()
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0)
        logger.info(f"Task {task_id}: Resume completed ({pages_processed} pages)")

    except Exception as e:
        logger.error(f"Task {task_id}: Resume error: {e}")
        import traceback
        logger.error(traceback.format_exc())
        with progress_lock:
            progress_tracker[task_id]["status"] = "error"
            progress_tracker[task_id]["error"] = str(e)
    finally:
        _release_keepalive()


def _ensure_keepalive():
    global _active_tasks, _keepalive_thread
    with _keepalive_lock:
        _active_tasks += 1
        if _keepalive_thread is None or not _keepalive_thread.is_alive():
            _keepalive_thread = threading.Thread(target=_keepalive_loop, daemon=True)
            _keepalive_thread.start()
            logger.info("Keepalive thread started")


def _release_keepalive():
    global _active_tasks
    with _keepalive_lock:
        _active_tasks -= 1
        if _active_tasks < 0:
            _active_tasks = 0


def _keepalive_loop():
    render_url = os.environ.get("RENDER_EXTERNAL_URL") or os.environ.get("RENDER_URL", "")
    local_url = f"http://localhost:{os.environ.get('PORT', 8080)}"

    while True:
        with _keepalive_lock:
            if _active_tasks <= 0:
                break
        success = False
        if render_url:
            try:
                requests.get(f"{render_url}/health", timeout=15)
                success = True
            except Exception:
                pass
        if not success:
            try:
                requests.get(f"{local_url}/health", timeout=5)
                success = True
            except Exception:
                pass
        if not success:
            logger.warning("Keepalive: all endpoints unreachable")
        time.sleep(60)  # Ping every 60s instead of 120s for faster response


def _ocr_keepalive_ping():
    """Secondary keepalive: called from within the OCR thread itself.
    Pings localhost every 60s so Render doesn't sleep even if the main
    keepalive thread is delayed or dead."""
    try:
        local_url = f"http://localhost:{os.environ.get('PORT', 8080)}"
        requests.get(f"{local_url}/health", timeout=5)
    except Exception:
        pass


# =====================================================================
# 24x7 Uptime keepalive - keeps the app awake at all times
# Render's free tier sleeps the instance after ~15 min of inactivity.
# A dedicated background thread pings our own /health endpoint on a
# regular interval (4 min) in both local (dev) and Render (prod) modes
# so the service stays online around the clock, whether or not OCR
# tasks are currently running.
# =====================================================================
_uptime_interval = int(os.environ.get("UPTIME_KEEPALIVE_INTERVAL", "240"))  # seconds
_uptime_thread = None
_uptime_lock = threading.Lock()


def _candidate_external_urls():
    """Return ordered candidate public URLs to ping so the app wakes itself up.

    Tries, in order:
      1. An explicitly configured UPTIME_KEEPALIVE_URL (most reliable if set).
      2. Render-provided env URLs (RENDER_EXTERNAL_URL / RENDER_URL / RANDOM_ENV_URL).
      3. The known custom domain(s) for this app (app.meowocr.work.gd and root).
      4. A .onrender.com URL derived from RENDER_SERVICE_NAME.
    """
    cands = []
    explicit = (os.environ.get("UPTIME_KEEPALIVE_URL") or "").strip().rstrip("/")
    if explicit and explicit not in cands:
        cands.append(explicit)
    for var in ("RENDER_EXTERNAL_URL", "RENDER_URL", "RANDOM_ENV_URL"):
        val = (os.environ.get(var) or "").strip().rstrip("/")
        if val and val not in cands:
            cands.append(val)
    # Known custom domain(s) for this deployment (www is the confirmed live URL)
    for host in ("www.meowocr.work.gd", "meowocr.work.gd"):
        url = f"https://{host}"
        if url not in cands:
            cands.append(url)
    # Derive default from service name (Render sets RENDER_SERVICE_NAME)
    svc = (os.environ.get("RENDER_SERVICE_NAME") or "").strip().lower()
    if svc:
        default = f"https://{svc}.onrender.com"
        if default not in cands:
            cands.append(default)
    return cands


def _uptime_keepalive_loop():
    port = os.environ.get("PORT", "8080")
    local_http = f"http://localhost:{port}/healthz"
    local_https_guess = f"https://localhost:{port}/healthz"
    external_urls = _candidate_external_urls()
    logger.info(
        "Uptime keepalive loop started (interval=%ss, external=%s)",
        _uptime_interval,
        external_urls or "none (localhost only)",
    )
    _last_err = ""
    loops = 0
    while True:
        success = False
        for base in external_urls:
            url = f"{base}/healthz"
            try:
                requests.get(url, timeout=15, verify=True)
                success = True
                if url != _last_err:
                    logger.debug("Uptime keepalive OK: %s", url)
                _last_err = ""
                break
            except Exception as exc:  # noqa: BLE001
                _last_err = f"{url} -> {exc}"
                continue
        if not success and external_urls:
            # Last-resort local checks keep local/dev instances alive too.
            for url in (local_http, local_https_guess):
                try:
                    requests.get(url, timeout=5, verify=False)
                    success = True
                    break
                except Exception:
                    continue
        if not success:
            logger.warning("Uptime keepalive: no reachable endpoint (%s)", _last_err)
        # Guard against a pathological zero/negative interval
        interval = _uptime_interval if _uptime_interval > 20 else 240
        loops += 1
        # Periodically surface a progress heartbeat so it's visible in logs
        if loops % (24 * 60 * 60 // interval) == 0:  # ~once per day
            logger.info("Uptime keepalive heartbeat: still running (interval=%ss)", interval)
        time.sleep(interval)


def start_uptime_keepalive():
    """Start the 24x7 keepalive background thread (idempotent).
    Called at module import so it runs under both `python app.py`
    (local dev) and gunicorn (Render), keeping the app awake."""
    global _uptime_thread
    if os.environ.get("DISABLE_UPTIME_KEEPALIVE", "").lower() in ("1", "true", "yes"):
        return
    with _uptime_lock:
        if _uptime_thread is None or not _uptime_thread.is_alive():
            _uptime_thread = threading.Thread(target=_uptime_keepalive_loop, daemon=True)
            _uptime_thread.start()


# Start the 24x7 keepalive as soon as the module loads.
start_uptime_keepalive()


def handle_translate_post():
    """Handle POST request for translation."""
    if 'file' not in request.files:
        flash('No file selected')
        return redirect(url_for('index'))

    file = request.files['file']

    if file.filename == '':
        flash('No file selected')
        return redirect(url_for('index'))

    ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
    if ext not in ('txt', 'md'):
        flash('Please upload a .txt or .md file for translation')
        return redirect(url_for('index'))

    source_lang = request.form.get('source_language', 'auto')
    target_lang = request.form.get('target_language', 'en')

    source_lang = TRANSLATOR_LANG_MAP.get(source_lang, source_lang)
    target_lang = TRANSLATOR_LANG_MAP.get(target_lang, target_lang)

    if not target_lang:
        flash('Please select a target language')
        return redirect(url_for('index'))

    filename = secure_filename(file.filename)
    temp_dir = tempfile.mkdtemp()
    file_path = os.path.join(temp_dir, filename)
    file.save(file_path)

    task_id = str(uuid.uuid4())
    with progress_lock:
        progress_tracker[task_id] = {
            "current_chunk": 0,
            "status": "starting",
            "error": None,
            "filename": filename,
            "temp_dir": temp_dir,
            "file_path": file_path,
            "percentage": 0,
            "total_chunks": None,
            "total_chars": None,
            "source_lang": source_lang,
            "target_lang": target_lang,
            "file_type": "translation",
            "created_at": time.time(),
            "cancelled": False,
            "translating": True
        }
    _save_progress(True)

    thread = threading.Thread(
        target=translate_file_background,
        args=(task_id, file_path, filename, temp_dir, source_lang, target_lang)
    )
    thread.daemon = True
    thread.start()
    _ensure_keepalive()

    return render_template('processing.html', task_id=task_id)


def translate_text(text, target_lang, source_lang='auto', chunk_size=2000):
    """Translate text using Google Translate's unofficial API with retries."""
    from bs4 import BeautifulSoup
    import requests

    session = requests.Session()
    session.headers.update({
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
    })
    base_url = 'https://translate.google.com/m'

    def _translate(text, max_retries=3):
        for attempt in range(max_retries):
            try:
                params = {'q': text, 'sl': source_lang, 'tl': target_lang}
                resp = session.get(base_url, params=params, timeout=15)
                if resp.status_code == 429:
                    wait = 30 * (attempt + 1)
                    logger.warning(f"Translate 429, waiting {wait}s (attempt {attempt+1}/{max_retries})")
                    time.sleep(wait)
                    continue
                if resp.status_code != 200:
                    if attempt < max_retries - 1:
                        wait = 5 * (2 ** attempt)
                        logger.warning(f"Translate status={resp.status_code}, retrying in {wait}s (attempt {attempt+1}/{max_retries})")
                        time.sleep(wait)
                        continue
                    raise Exception(f"HTTP {resp.status_code}")
                soup = BeautifulSoup(resp.text, 'html.parser')
                result = soup.find('div', class_='t0')
                if not result:
                    result = soup.find('div', class_='result-container')
                if result:
                    translated = result.get_text(strip=True)
                    if translated:
                        return translated
                if attempt < max_retries - 1:
                    time.sleep(5 * (attempt + 1))
                else:
                    raise Exception("No translation found in response")
            except requests.exceptions.Timeout:
                if attempt < max_retries - 1:
                    logger.warning(f"Translate timeout, retrying in {5*(attempt+1)}s (attempt {attempt+1}/{max_retries})")
                    time.sleep(5 * (attempt + 1))
                else:
                    raise
            except requests.exceptions.ConnectionError:
                if attempt < max_retries - 1:
                    logger.warning(f"Translate connection error, retrying in {5*(attempt+1)}s (attempt {attempt+1}/{max_retries})")
                    time.sleep(5 * (attempt + 1))
                else:
                    raise
        raise Exception(f"All {max_retries} attempts failed")

    if len(text) <= chunk_size:
        return _translate(text)

    result_parts = []
    total_chars = len(text)
    start = 0
    chunk_num = 0
    while start < total_chars:
        end = min(start + chunk_size, total_chars)
        if end < total_chars:
            break_at = text.rfind('.', start, end)
            if break_at > start:
                end = break_at + 1
        subchunk = text[start:end]
        try:
            translated = _translate(subchunk)
            result_parts.append(translated)
        except Exception as e:
            logger.warning(f"Translate subchunk {chunk_num} failed after all retries: {e}")
            result_parts.append(subchunk)
        start = end
        chunk_num += 1
    return ''.join(result_parts)


def translate_file_background(task_id, file_path, filename, temp_dir, source_lang, target_lang):
    """Background thread to translate a text file and update progress."""
    try:
        with progress_lock:
            progress_tracker[task_id]["status"] = "translating"
            progress_tracker[task_id]["detected_language"] = source_lang

        output_filename = f"{os.path.splitext(filename)[0]}_translated_{target_lang}.txt"
        output_path = os.path.join(temp_dir, output_filename)

        with progress_lock:
            progress_tracker[task_id]["output_path"] = output_path
            progress_tracker[task_id]["output_filename"] = output_filename

        text = process_txt_file(file_path)
        total_chars = len(text)
        if total_chars == 0:
            with progress_lock:
                progress_tracker[task_id]["status"] = "error"
                progress_tracker[task_id]["error"] = "Input file is empty"
            _save_progress(True)
            return

        with progress_lock:
            progress_tracker[task_id]["total_chars"] = total_chars

        chunk_size = 4500
        chunks = []
        start = 0
        while start < total_chars:
            end = min(start + chunk_size, total_chars)
            if end < total_chars:
                break_at = text.rfind('.', start, end)
                if break_at > start:
                    end = break_at + 1
            chunks.append(text[start:end])
            start = end

        total_chunks = len(chunks)
        with progress_lock:
            progress_tracker[task_id]["total_chunks"] = total_chunks

        translated_parts = []
        for i, chunk in enumerate(chunks):
            if task_id:
                with progress_lock:
                    if progress_tracker[task_id].get("cancelled"):
                        logger.info(f"Task {task_id}: Translation cancelled")
                        progress_tracker[task_id]["status"] = "cancelled"
                        progress_tracker[task_id]["error"] = "Cancelled by user"
                        _save_progress(True)
                        return

                try:
                    translated = translate_text(chunk, target_lang)
                    logger.info(f"Task {task_id}: Chunk {i+1}/{total_chunks} OK (len={len(translated)})")
                    if translated == chunk:
                        logger.warning(f"Task {task_id}: Chunk {i+1}/{total_chunks} returned SAME, retrying...")
                        time.sleep(5)
                        translated = translate_text(chunk, target_lang)
                        if translated == chunk:
                            logger.warning(f"Task {task_id}: Chunk {i+1}/{total_chunks} giving up, using original")
                    translated_parts.append(translated)
                except Exception as e:
                    logger.warning(f"Task {task_id}: Chunk {i+1}/{total_chunks} failed: {e}")
                    translated_parts.append(chunk)

                time.sleep(5)

                pct = int(((i + 1) / total_chunks) * 100)
                with progress_lock:
                    progress_tracker[task_id]["current_chunk"] = i + 1
                    progress_tracker[task_id]["percentage"] = pct
                _save_progress()

        result_text = ''.join(translated_parts)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(result_text)

        _save_progress()

        try:
            mega_link = upload_to_mega(output_path, output_filename)
            with progress_lock:
                progress_tracker[task_id]["download_link"] = mega_link
                if mega_link:
                    progress_tracker[task_id]["mega_uploaded"] = True
                    progress_tracker[task_id]["mega_status"] = "uploaded"
                    logger.info(f"Task {task_id}: Mega upload success - {mega_link}")
                else:
                    progress_tracker[task_id]["mega_uploaded"] = False
                    progress_tracker[task_id]["mega_status"] = ""
                    logger.warning(f"Task {task_id}: Mega upload returned None")
        except Exception as mega_err:
            logger.error(f"Task {task_id}: Mega upload error: {mega_err}")
            with progress_lock:
                progress_tracker[task_id]["mega_uploaded"] = False
                progress_tracker[task_id]["mega_status"] = ""

        persist_output(task_id)

        with progress_lock:
            progress_tracker[task_id]["status"] = "completed"
            progress_tracker[task_id]["percentage"] = 100
            progress_tracker[task_id]["completed_at"] = time.time()
            progress_tracker[task_id]["source_lang"] = source_lang
            progress_tracker[task_id]["target_lang"] = target_lang
        _save_progress(True)
        logger.info(f"Task {task_id}: Translation completed ({total_chunks} chunks, {total_chars} chars)")

    except Exception as e:
        logger.error(f"Task {task_id}: Translation error - {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        with progress_lock:
            progress_tracker[task_id]["status"] = "error"
            progress_tracker[task_id]["error"] = str(e)
        _save_progress(True)
    finally:
        _release_keepalive()


@app.route('/', methods=['GET', 'POST'])
@app.route('/translate', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Upload rate limit (uploads are the expensive path — blunt abuse)
        _rl = _check_rate(6, 60)  # max ~6 uploads per minute per IP
        if _rl is not None:
            return _rl
        # Check if translation request
        is_translate = request.path == '/translate'
        if is_translate or request.form.get('mode') == 'translate':
            return handle_translate_post()

        if 'file' not in request.files:
            flash('No file selected')
            return redirect(request.url)

        file = request.files['file']

        if file.filename == '':
            flash('No file selected')
            return redirect(request.url)

        if not allowed_file(file.filename):
            flash('Unsupported file type')
            return redirect(request.url)
        
        # Detect file type
        file_type = get_file_type(file.filename)
        if not file_type:
            flash('Unsupported file type')
            return redirect(request.url)

        # Enforce free-anonymous conversion limit (skip for logged-in users)
        if not session.get("logged_in") and not session.get("uid"):
            used = 0
            try:
                used = int(request.cookies.get(_COOKIE_COUNTER, "0") or "0")
            except (TypeError, ValueError):
                used = 0
            if used >= FREE_DOCS_WITHOUT_LOGIN and FIREBASE_ENABLED:
                flash(f"You have used your {FREE_DOCS_WITHOUT_LOGIN} free anonymous scan. Sign in with Google to convert unlimited documents - free!")
                return redirect(request.url)
        
        # Get language selection (only used for PDF and image)
        selected_lang = request.form.get('language', 'auto')
        logger.info(f"Language selection: {selected_lang}")
        
        # Get page range (only for PDF)
        page_range = request.form.get('page_range', 'all')
        start_page = 1
        end_page = None
        
        if page_range == 'custom':
            try:
                start_page = int(request.form.get('start_page', 1))
                end_page_str = request.form.get('end_page', '').strip()
                if end_page_str:
                    end_page = int(end_page_str)
                logger.info(f"Custom page range: {start_page} to {end_page}")
            except ValueError:
                logger.warning("Invalid page range, using defaults")
                start_page = 1
                end_page = None

        # Get auto-delete preference (days) - default 2 days
        try:
            auto_delete_days = int(request.form.get('auto_delete', '2'))
        except (TypeError, ValueError):
            auto_delete_days = AUTO_DELETE_DAYS

        # Save uploaded file to temp location
        filename = secure_filename(file.filename)
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, filename)
        file.save(file_path)
        
        # Verify file was saved
        if not os.path.exists(file_path):
            logger.error(f"Failed to save file to {file_path}")
            flash('Failed to save uploaded file')
            return redirect(request.url)
        
        logger.info(f"File saved to {file_path}, size: {os.path.getsize(file_path)} bytes")
        
        # Create a task ID for progress tracking
        task_id = str(uuid.uuid4())
        logger.info(f"Created task {task_id} for file {filename}")
        
        # Initialize progress tracker
        with progress_lock:
            progress_tracker[task_id] = {
                "current_page": 0,
                "status": "starting",
                "result_path": None,
                "error": None,
                "filename": filename,
                "temp_dir": temp_dir,
                "file_path": file_path,
                "total_pages": None,
                "percentage": 0,
                "detected_language": selected_lang if selected_lang != 'auto' else None,
                "selected_lang": selected_lang,
                "start_page": start_page,
                "end_page": end_page,
                "created_at": time.time(),
                "file_type": file_type,
                "cancelled": False,
                "auto_delete_days": auto_delete_days
            }
        _save_progress(True)
        
        # Start background processing thread
        logger.info(f"Starting background thread for task {task_id}, type: {file_type}")
        thread = threading.Thread(
            target=process_file_background,
            args=(task_id, file_path, filename, temp_dir, selected_lang, start_page, end_page, 200, file_type)
        )
        thread.daemon = True
        thread.start()
        logger.info(f"Background thread started for task {task_id}")
        _ensure_keepalive()
        
        # Render the processing page with task ID
        resp = make_response(render_template('processing.html', task_id=task_id))
        # Track anonymous conversions with a cookie (no limit for logged-in users)
        if not session.get("logged_in") and not session.get("uid"):
            used = 0
            try:
                used = int(request.cookies.get(_COOKIE_COUNTER, "0") or "0")
            except (TypeError, ValueError):
                used = 0
            resp.set_cookie(_COOKIE_COUNTER, str(used + 1), max_age=365 * 86400, httponly=True)
        return resp
    
    return render_template('index.html')

@app.route('/progress/<task_id>')
def get_progress(task_id):
    """Return progress status as JSON"""
    try:
        with progress_lock:
            if task_id not in progress_tracker:
                logger.warning(f"Progress check: task {task_id} not found in tracker")
                return jsonify({"status": "not_found"}), 404

            task = progress_tracker[task_id]
            
            # Check if this is a translation task
            is_translation = task.get("translating") or task.get("file_type") == "translation"
            
            if is_translation:
                source = task.get("source_lang", "auto")
                target = task.get("target_lang", "en")
                source_display = TRANSLATOR_TARGET_LANGS.get(source, source)
                target_display = TRANSLATOR_TARGET_LANGS.get(target, target)
                return jsonify({
                    "status": task["status"],
                    "percentage": task.get("percentage", 0),
                    "error": task["error"],
                    "filename": task["filename"],
                    "file_type": "translation",
                    "source_lang": source,
                    "source_lang_name": source_display,
                    "target_lang": target,
                    "target_lang_name": target_display,
                    "current_chunk": task.get("current_chunk", 0),
                    "total_chunks": task.get("total_chunks"),
                    "total_chars": task.get("total_chars"),
                    "translating": True
                })
            
            # Get language code - use selected language if available
            lang_code = task.get("detected_language")
            if not lang_code or lang_code == 'auto':
                lang_code = DEFAULT_LANG
            
            lang_name = LANG_CODE_TO_NAME.get(lang_code, lang_code)
            
            # Compute ETA
            eta = None
            if task["status"] in ("processing", "resuming"):
                page_times = task.get("page_times")
                total = task.get("total_pages")
                current_page = task.get("current_page", 0)
                start_page_offset = task.get("processing_start_page", 1)
                if page_times and total and current_page > 0:
                    avg = sum(page_times) / len(page_times)
                    pages_done = max(0, current_page - start_page_offset + 1)
                    remaining = max(0, total - pages_done)
                    if remaining > 0:
                        eta = int(avg * remaining)
            
            return jsonify({
                "status": task["status"],
                "current_page": task["current_page"],
                "total_pages": task.get("total_pages"),
                "pdf_total_pages": task.get("pdf_total_pages"),
                "percentage": task.get("percentage", 0),
                "error": task["error"],
                "filename": task["filename"],
                "detected_language": lang_code,
                "language_name": lang_name,
                "file_type": task.get("file_type", "pdf"),
                "eta": eta
            })
    except Exception as e:
        logger.error(f"Progress endpoint error for {task_id}: {e}")
        return jsonify({"status": "error", "error": str(e)}), 500


@app.route('/check/<task_id>')
def check_progress_page(task_id):
    """Styled HTML progress page for checking task status."""
    return render_template('progress_check.html', task_id=task_id)


@app.route('/download/<task_id>')
def download_result(task_id):
    """Download OCR result file"""
    with progress_lock:
        if task_id not in progress_tracker:
            flash('Task not found')
            return redirect(url_for('index'))

        task = progress_tracker[task_id]

        if task["status"] != "completed":
            flash('Processing not completed')
            return redirect(url_for('index'))

        if not task["output_path"] or not os.path.exists(task["output_path"]):
            if task.get("download_link"):
                with progress_lock:
                    progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1
                return redirect(task["download_link"])
            if task.get("mega_node_info"):
                with progress_lock:
                    progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1
                return redirect(url_for('get_mega_link', task_id=task_id))
            flash('Result file not found')
            return redirect(url_for('index'))

        output_path = task["output_path"]
        output_filename = task["output_filename"]

    # Send file as download
    fmt = request.args.get("format", "txt")
    try:
        if fmt == "docx":
            with open(output_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            base = os.path.splitext(output_filename)[0]
            pc = base.rsplit("_ocr", 1)[0] if base.endswith("_ocr") else base
            docx_name = f"{pc}_ocr.docx"
            buf = _build_docx(text)
            response = send_file(
                buf,
                as_attachment=True,
                download_name=docx_name,
                mimetype=_docx_mimetype()
            )
        else:
            response = send_file(
                output_path,
                as_attachment=True,
                download_name=output_filename,
                mimetype='text/plain'
            )
    except Exception as e:
        logger.error(f"Download error for {task_id}: {e}")
        flash('Error sending file')
        return redirect(url_for('index'))

    # Track download count
    with progress_lock:
        progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1

    return response


@app.route('/track-download/<task_id>', methods=['POST'])
def track_download(task_id):
    """Increment download count for cloud (Mega) link clicks"""
    with progress_lock:
        if task_id in progress_tracker:
            progress_tracker[task_id]["download_count"] = progress_tracker[task_id].get("download_count", 0) + 1
            return jsonify({"ok": True}), 200
    return jsonify({"ok": False}), 404


def _build_docx(text, title="Meow OCR — Extracted Text"):
    """Return a BytesIO of a .docx built from extracted text."""
    import io as _io
    from docx import Document
    from docx.shared import Pt
    buf = _io.BytesIO()
    doc = Document()
    if title:
        h = doc.add_heading(title, level=1)
        for r in h.runs:
            r.font.size = Pt(18)
    for para in (text or "").splitlines():
        doc.add_paragraph(para or " ")
    doc.save(buf)
    buf.seek(0)
    return buf


def _docx_mimetype():
    return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def _read_task_text(task_id):
    """Return the extracted text of a completed task (from disk or Mega), or None."""
    with progress_lock:
        if task_id not in progress_tracker:
            return None
        task = progress_tracker[task_id]
        if task["status"] != "completed":
            return None
        output_path = task.get("output_path") or ""
        filename = task.get("output_filename", "output.txt")
    if output_path and os.path.exists(output_path):
        try:
            with open(output_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Result read error for {task_id}: {e}")
    # Try Mega fallback
    if os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD"):
        try:
            from mega import Mega
            import tempfile as _tf
            m = Mega().login(os.environ.get("MEGA_EMAIL"), os.environ.get("MEGA_PWD"))
            node = mega_call(m, "find", f"ocr-outputs/{filename}", timeout=30)
            if node:
                with _tf.TemporaryDirectory() as td:
                    mega_call(m, "download", node, dest_path=td, timeout=120)
                    local = os.path.join(td, os.path.basename(filename))
                    if os.path.exists(local):
                        with open(local, "r", encoding="utf-8", errors="replace") as f:
                            return f.read()
        except Exception as e:
            logger.error(f"Result Mega read error for {task_id}: {e}")
    return None


@app.route('/result/<task_id>')
def result_page(task_id):
    """Show extracted text with copy-to-clipboard and inline edit."""
    text = _read_task_text(task_id)
    if text is None:
        flash('Result not found or not completed yet')
        return redirect(url_for('index'))
    with progress_lock:
        filename = progress_tracker[task_id].get("output_filename", "output.txt")
        display_name = progress_tracker[task_id].get("original_filename", filename) if progress_tracker[task_id].get("original_filename") else filename
        low_conf_words = progress_tracker[task_id].get("low_conf_words", [])
    return render_template('result.html', task_id=task_id, text=text, filename=display_name, low_conf_words=low_conf_words)


@app.route('/refresh-cloud', methods=['POST'])
def refresh_cloud():
    """Re-scan Mega ocr-outputs and restore any cloud files not yet in My Downloads."""
    if not (os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD")):
        flash('Mega cloud storage is not configured (MEGA_EMAIL/MEGA_PWD env vars missing).', 'error')
        return redirect(url_for('downloads_page'))
    try:
        rebuild_completed_from_mega()
        if _mega_restore_error:
            flash(f'Could not refresh cloud files: {_mega_restore_error}', 'error')
        else:
            flash('Cloud storage scan complete. Your cloud files are updated.', 'success')
    except Exception as e:
        flash(f'Cloud refresh failed: {e}', 'error')
    return redirect(url_for('downloads_page'))


@app.route('/api/refresh-cloud', methods=['POST'])
def refresh_cloud_api():
    """JSON endpoint to trigger an on-demand Mega cloud re-scan."""
    if not (os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD")):
        return jsonify({"ok": False, "error": "Mega cloud storage is not configured"}), 400
    try:
        rebuild_completed_from_mega()
        if _mega_restore_error:
            return jsonify({"ok": False, "error": _mega_restore_error}), 500
        return jsonify({"ok": True, "message": "Cloud storage scan complete"}), 200
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500


@app.route('/mega-link/<task_id>', methods=['GET', 'POST'])
def get_mega_link(task_id):
    """Generate and cache a Mega download link for a task on demand."""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return jsonify({"error": "Not found"}), 404
        if task.get("download_link"):
            return jsonify({"link": task["download_link"]}), 200
        output_filename = task.get("output_filename")
        if not output_filename:
            return jsonify({"error": "Output filename not known"}), 404
    try:
        from mega import Mega
        email = os.environ.get("MEGA_EMAIL")
        password = os.environ.get("MEGA_PWD")
        if not email or not password:
            return jsonify({"error": "Mega not configured"}), 500
        m = Mega().login(email, password)
        found = mega_call(m, "find", f"ocr-outputs/{output_filename}", timeout=30)
        if not found:
            return jsonify({"error": "File not found in Mega storage"}), 404
        node = found[0] if isinstance(found, list) else found
        link = mega_call(m, "get_link", node, timeout=30)
        if not link:
            return jsonify({"error": "Failed to get link"}), 500
        with progress_lock:
            task = progress_tracker.get(task_id)
            if task:
                task["download_link"] = link
        return jsonify({"link": link}), 200
    except Exception as e:
        logger.error(f"Failed to get Mega link for {task_id}: {e}")
        return jsonify({"error": str(e)}), 500


@app.route('/share-link/<task_id>')
def share_link(task_id):
    """Return a shareable download link for a completed task"""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task or task.get("status") != "completed":
            return jsonify({"error": "Not available"}), 404
        link = task.get("download_link")
        output_filename = task.get("output_filename")
    if link:
        return jsonify({"link": link})
    if output_filename:
        try:
            from mega import Mega
            email = os.environ.get("MEGA_EMAIL")
            password = os.environ.get("MEGA_PWD")
            if email and password:
                m = Mega().login(email, password)
                found = mega_call(m, "find", f"ocr-outputs/{output_filename}", timeout=30)
                if found:
                    node = found[0] if isinstance(found, list) else found
                    link = mega_call(m, "get_link", node, timeout=30)
                    if link:
                        with progress_lock:
                            t = progress_tracker.get(task_id)
                            if t:
                                t["download_link"] = link
                        return jsonify({"link": link})
        except Exception as e:
            logger.error(f"Share link Mega error for {task_id}: {e}")
    return jsonify({"link": url_for('download_result', task_id=task_id, _external=True)})


@app.route('/cancel/<task_id>', methods=['POST'])
def cancel_task(task_id):
    """Cancel a running OCR task"""
    with progress_lock:
        task = progress_tracker.get(task_id)
        if not task:
            return jsonify({"error": "Task not found"}), 404
        if task["status"] not in ("processing", "resuming", "starting", "detecting_language", "getting_page_count", "translating"):
            return jsonify({"error": "Task is not running"}), 400
        task["cancelled"] = True
        task["status"] = "cancelling"
    _save_progress(True)
    logger.info(f"Task {task_id}: Cancel requested")
    return jsonify({"status": "cancelling"}), 200


@app.route('/retry/<task_id>', methods=['POST'])
def retry_task(task_id):
    """Retry a failed OCR task"""
    with progress_lock:
        old = progress_tracker.get(task_id)
        if not old:
            return jsonify({"error": "Task not found"}), 404
        if old.get("status") not in ("error", "interrupted", "completed_with_error"):
            return jsonify({"error": "Only failed or interrupted tasks can be retried"}), 400

        filename = old.get("filename", "input.pdf")
        file_type = old.get("file_type", "pdf")
        temp_dir = old.get("temp_dir")
        file_path = old.get("file_path")
        selected_lang = old.get("selected_lang", "auto")
        last_ckpt = old.get("last_checkpoint_page", 0)
        current_page = old.get("current_page", 0)
        end_page = old.get("end_page")
        old_total = old.get("total_pages")
        mega_node = old.get("mega_original_handle")
        detected_lang = old.get("detected_language")

    # Everything after here is outside the lock

    # Try to locate the original file
    src = file_path or (os.path.join(temp_dir, filename) if temp_dir else None)
    needs_new_temp = False

    if not src or not os.path.exists(src):
        if mega_node:
            new_temp = tempfile.mkdtemp()
            src = os.path.join(new_temp, filename)
            try:
                m = init_mega()
                mega_call(m, "download", mega_node, dest_path=src, timeout=120)
                needs_new_temp = True
                logger.info(f"Task {task_id}: Downloaded original from Mega for retry")
            except Exception as e:
                shutil.rmtree(new_temp, ignore_errors=True)
                return jsonify({"error": f"Mega download failed: {e}"}), 400
        else:
            return jsonify({"error": "Original file not found and no Mega backup"}), 400

    # Create new task
    new_id = str(uuid.uuid4())
    new_temp_dir = tempfile.mkdtemp() if needs_new_temp else temp_dir
    if needs_new_temp:
        shutil.copy2(src, os.path.join(new_temp_dir, filename))
        src = os.path.join(new_temp_dir, filename)
    else:
        src = file_path

    # Resume from last known page, not from page 1
    resume_start = (last_ckpt or current_page or 0) + 1
    if resume_start < 1:
        resume_start = 1

    with progress_lock:
        progress_tracker[new_id] = {
            "current_page": resume_start - 1,
            "status": "starting",
            "result_path": None,
            "error": None,
            "filename": filename,
            "temp_dir": new_temp_dir,
            "file_path": src,
            "total_pages": old_total,
            "percentage": 0,
            "detected_language": detected_lang,
            "selected_lang": selected_lang,
            "start_page": resume_start,
            "end_page": end_page,
            "created_at": time.time(),
            "file_type": file_type,
            "cancelled": False,
            "retry_of": task_id
        }
    _save_progress(True)

    thread = threading.Thread(
        target=process_file_background,
        args=(new_id, src, filename, new_temp_dir, selected_lang, resume_start, end_page, 200, file_type)
    )
    thread.daemon = True
    thread.start()
    _ensure_keepalive()

    return jsonify({"task_id": new_id}), 200


@app.route('/ads.txt')
def ads_txt():
    """Advertisement authority file required by ad networks.
    Replace the example line below with the ads.txt content from your ad network."""
    content = (
        "# Replace this with the ads.txt content provided by your ad network\n"
        "# Example: google.com, pub-0000000000000000, DIRECT, f08c47fec0942fa0\n"
    )
    return app.response_class(content, mimetype='text/plain')


# Canonical public address of the site. Keep in sync with Render custom domain.
SITE_URL = os.environ.get("SITE_URL", "https://www.meowocr.work.gd").rstrip("/")


@app.route('/sitemap.xml')
def sitemap():
    """XML sitemap so Google can discover and index every static page."""
    pages = ["", "/how-to-use", "/privacy", "/terms", "/downloads", "/about", "/tamil-ocr", "/hindi-ocr", "/english-ocr"]
    today = "2026-08-30"
    urls = "".join(
        f"  <url>\n"
        f"    <loc>{SITE_URL}{p}</loc>\n"
        f"    <lastmod>{today}</lastmod>\n"
        f"    <changefreq>weekly</changefreq>\n"
        f"    <priority>{'1.0' if p == '' else '0.7'}</priority>\n"
        f"  </url>\n"
        for p in pages
    )
    xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">\n'
        f"{urls}"
        '</urlset>\n'
    )
    return app.response_class(xml, mimetype='application/xml')


@app.route('/robots.txt')
def robots():
    """robots.txt allowing all crawlers and pointing to the sitemap."""
    content = (
        "User-agent: *\n"
        "Allow: /\n"
        f"Sitemap: {SITE_URL}/sitemap.xml\n"
    )
    return app.response_class(content, mimetype='text/plain')


@app.route('/google77bfad89af69530a.html')
def google_verify():
    """Google Search Console ownership verification file (exact byte content)."""
    content = "google-site-verification: google77bfad89af69530a.html\n"
    return app.response_class(content, mimetype='text/plain')


@app.route('/healthz')
def healthz():
    """Lightweight liveness endpoint used by the self-keepalive thread.
    Does NOT spawn a Tesseract subprocess, so it always responds in a few
    ms. External uptime monitors (UptimeRobot, cron-job.org, etc.) should
    ping /healthz for the cheapest possible check."""
    return {'status': 'ok'}, 200


@app.route('/health')
def health():
    """Health check endpoint for Render"""
    try:
        version = pytesseract.get_tesseract_version()
        # Check if Tamil and English packs are available
        languages = pytesseract.get_languages()
        return {
            'status': 'healthy',
            'tesseract_version': str(version),
            'languages_available': languages
        }, 200
    except Exception as e:
        return {'status': 'unhealthy', 'error': str(e)}, 500


@app.errorhandler(413)
def too_large(e):
    flash('File too large. Please upload a smaller file.')
    return redirect(url_for('index'))


@app.errorhandler(500)
def internal_error(e):
    import traceback
    try:
        with open("oc_error.log", "a") as f:
            f.write("\n===== 500 %s =====\n" % request.path)
            traceback.print_exc(file=f)
    except Exception:
        logger.error("500 handler failed to write error log", exc_info=True)
    return "Something went wrong on our side. Please go back and try again.", 500


@app.route('/clear-downloads', methods=['POST'])
def clear_downloads():
    data = request.get_json(silent=True) or {}
    task_ids = data.get("task_ids")
    cleared = 0
    with progress_lock:
        if task_ids:
            for tid in task_ids:
                if tid in progress_tracker:
                    task = progress_tracker[tid]
                    op = task.get("output_path")
                    if op and os.path.exists(op):
                        try: os.remove(op)
                        except: pass
                    td = task.get("temp_dir")
                    if td:
                        shutil.rmtree(td, ignore_errors=True)
                    del progress_tracker[tid]
                    cleared += 1
        else:
            for tid in list(progress_tracker.keys()):
                task = progress_tracker[tid]
                op = task.get("output_path")
                if op and os.path.exists(op):
                    try: os.remove(op)
                    except: pass
                td = task.get("temp_dir")
                if td:
                    shutil.rmtree(td, ignore_errors=True)
                del progress_tracker[tid]
                cleared += 1
    if cleared:
        _save_progress(True)
    return jsonify({"cleared": cleared}), 200


@app.route('/how-to-use')
def how_to_use():
    """Detailed user instructions page."""
    return render_template('how_to_use.html')


@app.route('/privacy')
def privacy_page():
    """Privacy Policy page."""
    return render_template('privacy.html')


@app.route('/terms')
def terms_page():
    """Terms of Service page."""
    return render_template('terms.html')


LANGUAGE_PAGES = {
    'tamil-ocr': {
        'route': 'tamil_ocr',
        'lang_code': 'tam',
        'name': 'Tamil',
        'emoji': '🇮🇳',
        'title': 'Tamil OCR - Extract Tamil Text from PDF & Images (Free Online)',
        'description': 'Free Tamil OCR online: extract editable Tamil text from PDF, JPG, PNG and scanned documents. No signup, no watermark, supports Tamil PDF to text conversion.',
        'h1': 'Tamil OCR - Convert Tamil PDF & Images to Text',
        'intro': 'Upload a Tamil PDF, image or scanned document and convert it into clean, editable Tamil text in seconds. Meow OCR uses Tesseract engine with automatic language detection, so handwritten and printed Tamil both extract accurately.',
        'features': ['Free Tamil PDF to text', 'No signup or card required', 'Supports JPG, PNG, TIFF & PDF', 'Results auto-delete in 2 days'],
    },
    'hindi-ocr': {
        'route': 'hindi_ocr',
        'lang_code': 'hin',
        'name': 'Hindi',
        'emoji': '🇮🇳',
        'title': 'Hindi OCR - Extract Hindi Text from PDF & Images (Free Online)',
        'description': 'Free Hindi OCR online: extract editable Hindi (Devanagari) text from PDF, images and scanned documents. No signup, supports Hindi PDF to text conversion.',
        'h1': 'Hindi OCR - Convert Hindi PDF & Images to Text',
        'intro': 'Turn Hindi (Devanagari) PDFs, images and scans into editable text in seconds. Meow OCR detects Hindi automatically and extracts clean, copyable text — completely free with no account needed.',
        'features': ['Free Hindi PDF to text', 'No signup or card required', 'Devanagari text extraction', 'Results auto-delete in 2 days'],
    },
    'english-ocr': {
        'route': 'english_ocr',
        'lang_code': 'eng',
        'name': 'English',
        'emoji': '🇬🇧',
        'title': 'PDF to Text Converter - Extract English Text from PDF & Images (Free)',
        'description': 'Free online PDF to text: extract editable English text from PDF, images and scanned documents. No signup, no watermark, fast and private.',
        'h1': 'English PDF to Text - Free Online OCR',
        'intro': 'Convert English PDFs, images and scans into clean, editable text in seconds. Fast, private and free — no signup, no watermark, and your files auto-delete after 2 days.',
        'features': ['Free English PDF to text', 'No signup or card required', 'Supports PDF, JPG, PNG & more', 'Results auto-delete in 2 days'],
    },
}


@app.route('/tamil-ocr')
def tamil_ocr():
    return render_template('language_page.html', data=LANGUAGE_PAGES['tamil-ocr'])


@app.route('/hindi-ocr')
def hindi_ocr():
    return render_template('language_page.html', data=LANGUAGE_PAGES['hindi-ocr'])


@app.route('/english-ocr')
def english_ocr():
    return render_template('language_page.html', data=LANGUAGE_PAGES['english-ocr'])


@app.route('/about')
def about():
    """About page — who runs Meow OCR."""
    return render_template('about.html')


def get_user():
    """Return the logged-in user dict, or None."""
    uid = session.get("uid")
    if not uid:
        return None
    return {
        "uid": uid,
        "name": session.get("name", ""),
        "email": session.get("email", ""),
        "photo": session.get("photo", ""),
    }


def verify_firebase_token(id_token):
    """Verify a Firebase ID token via the Firebase REST API (no SDK needed).
    Returns user dict on success, or None. Uses requests with caching handled
    by Google's servers; for production you can switch to firebase-admin."""
    if not FIREBASE_API_KEY:
        return None
    url = f"https://identitytoolkit.googleapis.com/v1/accounts:lookup?key={FIREBASE_API_KEY}"
    try:
        resp = requests.post(url, json={"idToken": id_token}, timeout=15)
        if resp.status_code != 200:
            logger.warning(f"Firebase lookup failed: {resp.status_code} {resp.text[:200]}")
            return None
        data = resp.json()
        users = data.get("users") or []
        if not users:
            return None
        u = users[0]
        return {
            "uid": u.get("localId", ""),
            "email": u.get("email", ""),
            "name": u.get("displayName", ""),
            "photo": u.get("photoUrl", ""),
        }
    except Exception as e:
        logger.error(f"Firebase token verification error: {e}")
        return None


@app.route('/api/auth', methods=['POST'])
@_rate_limit(20, 60)
def api_auth():
    """Verify Google (Firebase) ID token and start a server session."""
    data = request.get_json(silent=True) or {}
    id_token = data.get("idToken")
    if not id_token:
        return jsonify({"error": "Missing token"}), 400
    user = verify_firebase_token(id_token)
    if not user or not user.get("uid"):
        return jsonify({"error": "Authentication failed"}), 401
    session["uid"] = user["uid"]
    session["name"] = user.get("name", "")
    session["email"] = user.get("email", "")
    session["photo"] = user.get("photo", "")
    # Logged-in users get the free anonymous quota reset benefit (history persists server-side)
    session["logged_in"] = True
    return jsonify({"ok": True, "user": user}), 200


@app.route('/api/logout', methods=['POST'])
def api_logout():
    session.clear()
    return jsonify({"ok": True}), 200


@app.route('/api/me')
def api_me():
    user = get_user()
    return jsonify({
        "logged_in": bool(user),
        "user": user,
        "firebase_enabled": FIREBASE_ENABLED,
    })


@app.route('/downloads')
def downloads_page():
    all_tasks = []
    restoring = not _mega_restore_done.is_set()
    with progress_lock:
        for task_id, task in progress_tracker.items():
            status = task.get("status", "")

            # Compute ETA for in-progress tasks
            eta = None
            if status in ("processing", "resuming"):
                page_times = task.get("page_times")
                total = task.get("total_pages")
                current_page = task.get("current_page", 0)
                start_page_offset = task.get("processing_start_page", 1)
                if page_times and total and current_page > 0:
                    avg = sum(page_times) / len(page_times)
                    pages_done = max(0, current_page - start_page_offset + 1)
                    remaining = max(0, total - pages_done)
                    if remaining > 0:
                        eta = int(avg * remaining)

            is_translation = task.get("translating") or task.get("file_type") == "translation"
            info = {
                "task_id": task_id,
                "filename": task.get("output_filename", task.get("filename", "Unknown")),
                "download_link": task.get("download_link", ""),
                "language": task.get("detected_language", ""),
                "file_type": task.get("file_type", ""),
                "pages_processed": task.get("pages_processed", 0),
                "mega_uploaded": task.get("mega_uploaded", False),
                "mega_status": task.get("mega_status", ""),
                "completed_at": task.get("completed_at", 0),
                "created_at": task.get("created_at", 0),
                "download_count": task.get("download_count", 0),
                "status": status,
                "percentage": task.get("percentage", 0),
                "eta": eta,
                "last_checkpoint_page": task.get("last_checkpoint_page"),
                "translating": is_translation,
                "source_lang": task.get("source_lang", ""),
                "target_lang": task.get("target_lang", ""),
                "current_chunk": task.get("current_chunk", 0),
                "total_chunks": task.get("total_chunks")
            }
            all_tasks.append(info)
    all_tasks.reverse()
    return render_template(
        "downloads.html",
        downloads=all_tasks,
        restoring=restoring,
        cloud_restore_error=_mega_restore_error,
        mega_configured=bool(os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD")),
    )



# =====================================================================
# Groq AI — Meow Assistant (chatbot) + Handwritten Notes scanner.
# GROQ_API_KEY is read server-side only (Render env var) and is NEVER
# rendered to templates or static assets. If no key is set, the chatbot
# falls back to scripted FAQ replies and the handwriting endpoint returns
# a friendly error, so the site never breaks without a key.
# =====================================================================
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "").strip()
GROQ_CHAT_MODEL = os.environ.get("GROQ_CHAT_MODEL", "llama-3.1-8b-instant")
GROQ_VISION_MODEL = os.environ.get("GROQ_VISION_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct")
# Fallback vision models tried in order if the primary (or a configured one) is
# unavailable on this account (e.g. model decommissioned or not granted).
GROQ_VISION_FALLBACKS = [
    "meta-llama/llama-4-scout-17b-16e-instruct",
    "llama-4-scout-17b-16e-instruct",
    "qwen/qwen3.6-27b",
    "qwen/qwen3.8-27b",
    "llama-3.2-11b-vision-preview",
]
_GROQ_VISION_RETRY_4XX = 3  # how many times to retry a quirk transient 429/5xx before giving up
_GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"
_GROQ_MAX_MESSAGE = 500
_GROQ_MAX_IMAGE_BYTES = 25 * 1024 * 1024  # covers typical 15-20MB camera photos
_GROQ_MAX_IMAGE_PIXELS = 52_000_000  # ~7200x7200, lets 50MP camera photos pass
_GROQ_PREVIEW_SIZE = (1800, 1800)  # what the AI actually sees (payload stays <1MB)
_GROQ_DRAFT_THRESHOLD = 15_000_000  # big JPEGs decode at reduced scale to protect RAM
_GROQ_READ_CHUNK = 256 * 1024

_GROQ_SYSTEM_PROMPT = (
    "You are 'Meow Assistant', the friendly cat mascot of Meow OCR "
    "(https://www.meowocr.work.gd), a completely FREE online OCR tool "
    "created by developer Siva. You answer in a warm, playful but "
    "professional tone, using cat emojis sparingly (max one per message). "
    "Keep answers short (max ~100 words) and match the user's language "
    "(e.g. reply in Tamil or Hindi if they write in it). Answer only from "
    "these facts about Meow OCR:\n"
    "- Extracts editable text from PDF, images (JPG, PNG, TIFF), Word, "
    "Excel, PPT and documents, plus translate.\n"
    "- Supports 19+ languages incl. Tamil, Hindi, Telugu, Bengali, "
    "English, with auto-detection.\n"
    "- 100% free. First scan needs no signup; guests get one free scan, "
    "signing in with Google unlocks unlimited conversions.\n"
    "- No watermark, no credit card.\n"
    "- Files stored securely in the cloud, auto-deleted after 2 days by "
    "default; users can pick 1/7/30 days or keep forever and delete "
    "manually from My Downloads.\n"
    "- Handwritten notes: use the 'Handwritten Notes' AI scanner on the "
    "homepage to read photos of handwriting with AI.\n"
    "- Built by developer Siva.\n"
    "For anything outside these facts, honestly say you're not sure and "
    "suggest visiting the homepage help section."
)


def _faq_reply(text):
    """Keyword-matched scripted replies used when Groq is not configured,
    rate-limited, or unavailable — keeps the widget working at all times."""
    t = " " + text.lower().strip() + " "
    if any(k in t for k in ("hi ", "hello", "hey", "vanakkam", "good morning", "good evening")):
        return ("Hi! 🐱 I'm Meow Assistant. Ask me how Meow OCR works, whether it's "
                "free, which languages it supports, or how your files stay private.")
    if any(k in t for k in ("hand", "written", "handwriting", "note", "writing", "letra")):
        return ("Yes! Meow OCR can read handwriting 🐱 Use the 'Handwritten Notes' "
                "scanner on the homepage — upload a clear photo of your notes (or "
                "manuscript) and AI will transcribe them into editable text. "
                "Tamil, Hindi and English handwriting all work.")
    if any(k in t for k in ("free", "cost", "price", "pay", "money", "card", "subscription")):
        return ("It's 100% free 🐾 No payment, no trial, no credit card. Your first "
                "scan needs no signup at all, and signing in with Google unlocks "
                "unlimited conversions.")
    if any(k in t for k in ("language", "tamil", "hindi", "telugu", "supported")):
        return ("Meow OCR supports 19+ languages 🐱 including Tamil, Hindi, Telugu, "
                "Bengali, Malayalam, English, Arabic and more — with auto-detection "
                "so you usually don't need to pick manually. You can also translate "
                "text between 30+ languages.")
    if any(k in t for k in ("private", "privacy", "secure", "delete", "stored", "data", "safe")):
        return ("Your files are handled privately 🐱 They're stored securely in the "
                "cloud and automatically deleted after 2 days by default. You can "
                "choose 1/7/30 days or keep forever, and delete any download "
                "anytime from 'My Downloads'. No account is needed for basic use.")
    if any(k in t for k in ("how", "use", "upload", "works", "step", "convert", "start")):
        return ("Easy! 🐾 1) Go to the homepage. 2) Drag & drop your PDF, image or "
                "document (or tap the 'Handwritten Notes' card for handwriting). "
                "3) Hit 'Extract Text' — the engine reads it, then you can view, "
                "copy or download the .txt result. That's it.")
    if any(k in t for k in ("image", "jpg", "png", "photo", "picture", "pdf", "word", "excel", "format", "file type")):
        return ("You can upload PDFs, images (JPG, PNG, BMP, TIFF, GIF), Word, "
                "Excel, PPT, and more 🐾 For handwriting, use the 'Handwritten "
                "Notes' AI scanner instead — it's built for photos of written notes.")
    if any(k in t for k in ("contact", "siva", "who", "about", "developer", "created", "built", "email")):
        return ("Meow OCR was built by developer Siva 🐱 He keeps it free because "
                "he believes OCR should be simple. You can read more on the About "
                "page — there's a link at the bottom of every page.")
    if any(k in t for k in ("rule", "translate", "translation", "translate")):
        return ("Yes, Meow OCR can translate too! 🐾 Switch to the 'Translate' tab "
                "on the homepage, upload a .txt file, pick source and target "
                "languages, and download your translated file.")
    if any(k in t for k in ("error", "not working", "fail", "problem", "issue", "bug")):
        return ("Sorry about that! 🐱 Try a smaller file, use a clear high-quality "
                "scan, and pick the language manually for speedier results. If it "
                "still fails, your files are safe — wait a minute and try again.")
    return ("Great question! 🐱 I'm still a small cat assistant, and my long-term "
            "memory is short. Could you rephrase, or tap a quick question below? "
            "I'm best at questions about how Meow OCR works, free limits, "
            "languages, handwriting and privacy.")


@app.route("/api/chat", methods=["POST"])
@_rate_limit(20, 60)
def api_chat():
    """Meow Assistant chatbot. Accepts JSON {"message": "..."}, returns
    {"reply": "..."}. Calls Groq from the server; the API key never leaves
    the server. Falls back to _faq_reply when Groq is unavailable."""
    data = request.get_json(silent=True) or {}
    _m = data.get("message")
    raw = _m.strip() if isinstance(_m, str) else ""
    if not raw:
        return jsonify({"reply": "Say something, meow? 🐱"}), 200
    msg = raw[:_GROQ_MAX_MESSAGE]
    if not GROQ_API_KEY:
        return jsonify({"reply": _faq_reply(msg)}), 200
    try:
        resp = requests.post(
            _GROQ_URL,
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": GROQ_CHAT_MODEL,
                "messages": [
                    {"role": "system", "content": _GROQ_SYSTEM_PROMPT},
                    {"role": "user", "content": msg},
                ],
                "temperature": 0.4,
                "max_tokens": 500,
            },
            timeout=20,
        )
        if resp.status_code == 200:
            body = resp.json()
            reply = (body.get("choices") or [{}])[0].get("message", {}).get("content", "").strip()
            if reply:
                return jsonify({"reply": reply}), 200
        logger.warning("Groq chat non-200 status=%s", resp.status_code)
    except Exception:
        logger.error("Groq chat error, falling back to FAQ", exc_info=True)
    return jsonify({"reply": _faq_reply(msg)}), 200


def _clean_vision_output(text):
    """Strip any chain-of-thought/thinking preamble a vision model may prepend,
    and collapse accidental duplicate lines, returning just the transcribed text."""
    if not text:
        return text
    t = text.strip()
    # Some reasoning models prefix the answer with a "thinking" block. Keep only
    # the part after the closing "response" marker if one is present.
    marker = None
    for line in t.splitlines():
        low = (line or "").strip().lower()
        if low == "response" or low == "answer":
            marker = low
            break
    if marker is not None:
        parts = t.splitlines()
        for i, line in enumerate(parts):
            if (line or "").strip().lower() == marker:
                t = "\n".join(parts[i + 1:]).strip()
                break
    # If a raw "thinking" block still leads (no answer marker), cut everything
    # from "thinking" to the first blank-line+content transition.
    if t.lower().startswith("thinking") or "\nthinking\n" in t.lower() or t.lower().startswith("\nthinking"):
        t = re.sub(r"(?is)^\s*thinking\b.*?(\nresponse\b|\n\n)", "\n", t, count=1).strip()
        if t.lower().startswith("thinking"):
            t = ""
    # De-duplicate a classic "thinking echo" pattern where the model returns every
    # line twice (A,A,B,B,C,C). Only collapse when the whole block clearly shows
    # the alternating echo, so genuine repeated lines in handwritten poems/lyrics
    # are preserved.
    lines = t.splitlines()
    stripped = [ln.lstrip() for ln in lines]
    echo_ratio = sum(1 for i in range(len(stripped) - 1) if stripped[i] == stripped[i + 1]) / max(1, len(stripped) - 1)
    if echo_ratio >= 0.8:
        out = []
        for ln in lines:
            if out and out[-1].strip() == ln.strip():
                continue
            out.append(ln)
        return "\n".join(out).strip()
    return "\n".join(lines).strip()


@app.route("/api/handwrite", methods=["POST"])
@_rate_limit(10, 60)
def api_handwrite():
    """AI scanner for handwritten notes. Accepts an image file upload,
    sends it to the Groq vision model, returns {"text": "..."}. Runs fully
    server-side; nothing is stored and no key reaches the browser."""
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400
    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "No file selected."}), 400
    blob = b""
    while True:  # cap memory: reject oversized uploads while streaming, not after
        chunk = f.stream.read(_GROQ_READ_CHUNK)
        if not chunk:
            break
        blob += chunk
        if len(blob) > _GROQ_MAX_IMAGE_BYTES:
            return jsonify({"error": "Image is too large (max 25 MB)."}), 413
    if not blob:
        return jsonify({"error": "The file appears to be empty."}), 400
    try:
        img = Image.open(io.BytesIO(blob))
        # Reject decompression bombs before decoding (huge declared size in a tiny file)
        if (img.width or 0) * (img.height or 0) > _GROQ_MAX_IMAGE_PIXELS:
            return jsonify({"error": "Image resolution is too high (max ~7200 x 7200 pixels)."}), 413
        img.verify()  # cheap structural check; discards the decoded buffer
        img = Image.open(io.BytesIO(blob))
        w, h = img.size
        if (w or 0) * (h or 0) > _GROQ_DRAFT_THRESHOLD and img.format == "JPEG":
            img.draft(("RGB",), _GROQ_PREVIEW_SIZE)  # decode big JPEGs at reduced scale: low RAM spike
        img = img.convert("RGB")
        img.thumbnail(_GROQ_PREVIEW_SIZE)  # keep payload small + fast
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        data_url = "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()
    except Exception:
        logger.error("Handwrite image parse failed", exc_info=True)
        return jsonify({"error": "Could not read that image — please use JPG or PNG."}), 400
    if not GROQ_API_KEY:
        return jsonify({
            "error": "The AI handwriting engine isn't configured yet. Please try again later.",
            "ai_offline": True,
        }), 503
    prompt = (
        "This is a photo of HANDWRITTEN notes. Transcribe ALL handwritten or "
        "printed text exactly as written, preserving line breaks and paragraphs. "
        "If the handwriting is in Tamil, Hindi, or another language, transcribe "
        "it fully in that language. Do not add, correct or remove any words. "
        "Do not describe the image. If there is no readable handwriting, reply "
        "exactly: '(No readable text detected on this image)'."
    )
    # Build the candidate model list (primary first, then fallbacks), dedup in order.
    models = []
    for m in [GROQ_VISION_MODEL] + GROQ_VISION_FALLBACKS:
        m = (m or "").strip()
        if m and m not in models:
            models.append(m)
    last_status = None
    last_reason = None
    for model in models:
        try:
            payload = {
                "model": model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": data_url}},
                        ],
                    }
                ],
                "temperature": 0.2,
                "max_tokens": 2000,
            }
            # Qwen vision models default to a "thinking" mode that prefixes the
            # answer with chain-of-thought. Disable it for clean transcription.
            if "qwen" in (model or "").lower():
                payload["reasoning_effort"] = "none"
            resp = requests.post(
                _GROQ_URL,
                headers={
                    "Authorization": f"Bearer {GROQ_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=payload,
                timeout=60,
            )
            last_status = resp.status_code
            if resp.status_code == 200:
                body = resp.json()
                text = (body.get("choices") or [{}])[0].get("message", {}).get("content", "").strip()
                text = _clean_vision_output(text)
                if text:
                    return jsonify({"text": text}), 200
            last_reason = resp.text[:300]
            if resp.status_code == 400 and "decommissioned" in resp.text.lower():
                logger.warning("Vision model %s decommissioned, trying next", model)
                continue
            if resp.status_code in (401, 403):
                logger.warning("Vision model %s auth error HTTP %s: %s", model, resp.status_code, resp.text[:200])
                return jsonify({
                    "error": "The AI scanner has an invalid API key. Contact the site owner to fix the server setup.",
                    "ai_offline": True,
                }), 502
            if resp.status_code in (404, 400):
                logger.warning("Vision model %s unavailable HTTP %s: %s", model, resp.status_code, resp.text[:200])
                continue
            # other 4xx (e.g. 429 rate limit)
            if 400 <= resp.status_code < 500:
                logger.warning("Vision model %s HTTP %s: %s", model, resp.status_code, resp.text[:200])
                continue
            # 5xx server error: don't burn through fallbacks pointlessly
            logger.warning("Groq vision HTTP %s with model %s", resp.status_code, model)
        except Exception:
            logger.error("Groq vision error with model %s", model, exc_info=True)
            last_reason = "network error"
            continue
    if last_status in (401, 403):
        return jsonify({
            "error": "The AI scanner has an invalid API key. Contact the site owner to fix the server setup.",
            "ai_offline": True,
        }), 502
    if last_status and 400 <= last_status < 500:
        return jsonify({
            "error": "The AI scanner hit a snag (HTTP %s). Please try again in a moment." % last_status,
        }), 502
    return jsonify({"error": "The AI scanner is busy right now — please try again in a moment."}), 502


# Fast local restore (synchronous, <1s), Mega scan in background (fast with deferred links)
os.makedirs(OUTPUT_DIR, exist_ok=True)
saved = _load_progress()
if saved:
    for tid, data in saved.items():
        if data.get("status") not in ("completed", "error", "cancelled"):
            data["status"] = "interrupted"
        progress_tracker[tid] = data
    logger.info(f"Restored {len(saved)} persisted tasks from {PROGRESS_FILE}")

rebuild_completed_from_local()

with progress_lock:
    for tid, data in progress_tracker.items():
        if data.get("status") == "completed":
            op = data.get("output_path")
            if op and not os.path.exists(op):
                out_fn = data.get("output_filename")
                alt = os.path.join(OUTPUT_DIR, f"{tid}_{out_fn}") if out_fn else None
                if alt and os.path.exists(alt):
                    data["output_path"] = alt

if os.environ.get("MEGA_EMAIL") and os.environ.get("MEGA_PWD"):
    def _mega_background():
        try:
            rebuild_completed_from_mega()
            scan_and_resume_checkpoints()
        except Exception:
            logger.error("Mega background restore failed", exc_info=True)
        finally:
            _mega_restore_done.set()
    threading.Thread(target=_mega_background, daemon=True).start()
else:
    _mega_restore_done.set()

# Initial auto-delete pass on startup
try:
    cleanup_old_tasks()
except Exception:
    pass

# Periodic cleanup of old temp files (every hour)
def _cleanup_loop():
    while True:
        time.sleep(3600)
        try:
            cleanup_old_tasks()
        except Exception:
            pass

threading.Thread(target=_cleanup_loop, daemon=True).start()

if __name__ == '__main__':
    # Local testing only (Render uses Gunicorn)
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, threaded=True)
